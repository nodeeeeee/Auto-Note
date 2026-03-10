"""
Semantic Alignment Engine
Maps Whisper transcript segments → slide pages using dense RAG.

Pipeline:
  1. Extract text from each slide (PDF / PPTX / DOCX)
     - Includes speaker notes (PPTX) and image descriptions (OCR + Claude vision)
  2. Embed slides with sentence-transformers (all-mpnet-base-v2, GPU)
  3. Build a FAISS index for fast cosine-similarity lookup
  4. For each transcript segment query the index (optionally with a context
     window of ±CONTEXT_SEC seconds for richer matching signal)
  5. Apply Viterbi temporal smoothing so slides only advance forward
  6. Flag segments that don't match any slide well (off_slide)
  7. Collapse consecutive equal-slide segments into a compact timeline
  8. Save JSON to [course_id]/alignment/[stem].json

Usage:
  # align one caption↔slide pair
  python semantic_alignment.py \\
      --caption  85427/captions/CS3210\\ e-Lecture\\ on\\ Processes\\ and\\ Threads.json \\
      --slides   85427/materials/LectureNotes/L02-Processes-Threads.pdf

  # auto-discover all unaligned pairs in a course folder
  python semantic_alignment.py --course 85427
"""

from __future__ import annotations

import argparse
import io
import json
import sys
from pathlib import Path
from typing import NamedTuple

import faiss
import numpy as np

PROJECT_DIR = Path(__file__).parent

# ── Tunable knobs ─────────────────────────────────────────────────────────────

EMBED_MODEL   = "all-mpnet-base-v2"   # highest-quality general sentence model
CONTEXT_SEC   = 30.0                  # seconds of transcript to pool per query
BATCH_SIZE    = 64                    # embedding batch size

# Viterbi transition log-probabilities
STAY_LOGP     =  0.0    # free to stay on the same slide
FWD_LOGP_PER  = -0.02   # cost per slide advanced forward (nearly free)
BWD_LOGP_PER  = -1.5    # cost per slide stepped backward
                         # firm — allows genuine multi-minute review sections
                         # (e.g. slide 63 revisited) while suppressing
                         # single-segment noise flips

# Temporal position prior: at time t, add a Gaussian bonus centred on the
# expected slide position (t / duration) * n_slides.
# PRIOR_SIGMA controls the width in slide units; larger = softer guide.
PRIOR_SIGMA   = 8.0

# Off-slide detection: segments where the best raw cosine similarity (before
# the position prior) is below this threshold are flagged as off_slide.
OFF_SLIDE_THRESHOLD = 0.28

# Image captioning: pages / slides with fewer than this many words of text
# have their visual content described via OCR + Claude vision.
IMAGE_WORD_THRESHOLD = 20

# OpenAI vision model used for slide image description.
# gpt-4o-mini is used for cost efficiency; each call is only made once per
# slide because results are cached in {slide_file}.image_cache.json.
OPENAI_VISION_MODEL = "gpt-4o-mini"

# Sparse-slide enrichment: slides below this word count borrow text from
# content-rich neighbours before embedding.
SPARSE_THRESHOLD = 30   # words
NEIGHBOR_WORDS   = 60   # how many words to borrow from each neighbour


# ── OpenAI API key helper ─────────────────────────────────────────────────────

def _get_openai_key() -> str:
    import os
    key = os.environ.get("OPENAI_API_KEY", "")
    if not key:
        key_file = PROJECT_DIR / "openai_api.txt"
        if key_file.exists():
            key = key_file.read_text().strip()
    return key


# ── Image description (OCR + Claude vision) ───────────────────────────────────

_VISION_PROMPT = (
    "This is a slide from a university lecture. "
    "Describe ALL visible content in detail: every text label, concept name, "
    "diagram element, arrow, relationship, numbered step, code snippet, and "
    "technical term you can see. "
    "Focus on the academic topic — mention state names, algorithm steps, "
    "data structure relationships, or any specific terminology shown. "
    "Write a thorough plain-text description (no markdown), 3-5 sentences."
)


class ImageDescriber:
    """
    Extracts textual descriptions from slide images using:
      1. pytesseract OCR  — free, captures printed text in diagrams
      2. Claude vision API — powerful, understands academic diagrams and
                             technical content (used only on slides with
                             < IMAGE_WORD_THRESHOLD words; results cached)

    The Claude API is called at most once per slide per slide file, thanks to
    the image cache written to {slide_file}.image_cache.json.
    """

    def __init__(self):
        self._ocr_ok    = None   # None = unchecked, True/False after first call
        self._openai    = None   # openai.OpenAI client, lazy-loaded
        self._api_calls = 0      # count for cost reporting

    # ── lazy loaders ─────────────────────────────────────────────────────────

    def _ensure_ocr(self) -> bool:
        if self._ocr_ok is None:
            try:
                import pytesseract
                pytesseract.get_tesseract_version()
                self._ocr_ok = True
            except Exception:
                print("  [img] pytesseract not available — OCR skipped")
                self._ocr_ok = False
        return self._ocr_ok

    def _ensure_openai(self) -> bool:
        if self._openai is not None:
            return True
        key = _get_openai_key()
        if not key:
            print("  [img] No OpenAI API key — vision description skipped")
            return False
        try:
            from openai import OpenAI
            self._openai = OpenAI(api_key=key)
            return True
        except Exception as e:
            print(f"  [img] OpenAI client error: {e}")
            return False

    # ── per-image methods ────────────────────────────────────────────────────

    def _ocr(self, img) -> str:
        if not self._ensure_ocr():
            return ""
        try:
            import pytesseract
            return pytesseract.image_to_string(img).strip()
        except Exception:
            return ""

    def _openai_describe(self, img) -> str:
        """Send one slide image to GPT-4o-mini vision and return description."""
        if not self._ensure_openai():
            return ""
        import base64
        try:
            buf = io.BytesIO()
            img.save(buf, format="PNG")
            b64 = base64.standard_b64encode(buf.getvalue()).decode()

            resp = self._openai.chat.completions.create(
                model=OPENAI_VISION_MODEL,
                max_tokens=300,
                messages=[{
                    "role": "user",
                    "content": [
                        {"type": "image_url",
                         "image_url": {"url": f"data:image/png;base64,{b64}",
                                       "detail": "low"}},   # "low" = cheapest
                        {"type": "text", "text": _VISION_PROMPT},
                    ],
                }],
            )
            self._api_calls += 1
            return resp.choices[0].message.content.strip()
        except Exception as e:
            print(f"  [img] OpenAI vision error: {e}")
            return ""

    # ── public API ───────────────────────────────────────────────────────────

    def describe_slide_image(self, img) -> str:
        """
        Run OCR then Claude vision on a single slide image (already rendered
        as a full-page pixmap).  Returns combined description text.
        The caller is responsible for caching so this is never called twice
        for the same slide.
        """
        from PIL import Image as PILImage
        if not isinstance(img, PILImage.Image):
            return ""
        parts: list[str] = []
        ocr = self._ocr(img)
        if ocr:
            parts.append(ocr)
        vision = self._openai_describe(img)
        if vision:
            parts.append(vision)
        return " ".join(parts)

    @property
    def api_calls(self) -> int:
        return self._api_calls


# ── Image cache helpers ────────────────────────────────────────────────────────

def _load_image_cache(slide_path: Path) -> dict:
    cache_file = slide_path.parent / f"{slide_path.name}.image_cache.json"
    if cache_file.exists():
        with open(cache_file) as f:
            return json.load(f)
    return {}


def _save_image_cache(slide_path: Path, cache: dict) -> None:
    cache_file = slide_path.parent / f"{slide_path.name}.image_cache.json"
    with open(cache_file, "w") as f:
        json.dump(cache, f, indent=2)


# ── Slide text extraction ─────────────────────────────────────────────────────

class SlideText(NamedTuple):
    index: int          # 0-based slide/page index
    label: str          # short title (first non-empty line)
    text:  str          # full extracted text (incl. notes + image descriptions)


def extract_pdf(path: Path,
                describer: ImageDescriber | None = None,
                img_cache: dict | None = None) -> list[SlideText]:
    import fitz
    from PIL import Image as PILImage

    doc    = fitz.open(str(path))
    slides = []
    for i, page in enumerate(doc):
        text  = page.get_text().strip()
        label = next((ln.strip() for ln in text.splitlines() if ln.strip()),
                     f"Page {i+1}")

        # Image enrichment for sparse pages — render the full page as a
        # high-res pixmap so vector diagrams (state graphs, flowcharts, etc.)
        # are captured alongside any embedded bitmaps.
        if describer is not None and len(text.split()) < IMAGE_WORD_THRESHOLD:
            cache_key = f"page_{i}"
            if img_cache is not None and cache_key in img_cache:
                img_desc = img_cache[cache_key]
            else:
                pxmap = page.get_pixmap(matrix=fitz.Matrix(2, 2))
                pil   = PILImage.frombytes(
                    "RGB", [pxmap.width, pxmap.height], pxmap.samples)
                img_desc = describer.describe_slide_image(pil)
                if img_cache is not None:
                    img_cache[cache_key] = img_desc

            if img_desc:
                text = (text + "\n" + img_desc).strip()

        slides.append(SlideText(i, label[:80], text))
    doc.close()
    return slides


def extract_pptx(path: Path,
                 describer: ImageDescriber | None = None,
                 img_cache: dict | None = None) -> list[SlideText]:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from PIL import Image as PILImage

    prs    = Presentation(str(path))
    slides = []
    for i, slide in enumerate(prs.slides):
        parts: list[str] = []

        # Shape text
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = para.text.strip()
                    if line:
                        parts.append(line)

        # Speaker notes
        try:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                parts.append(notes_text)
        except Exception:
            pass

        text  = "\n".join(parts)
        label = parts[0][:80] if parts else f"Slide {i+1}"

        # Image enrichment for sparse slides: use the first picture shape.
        # For PPTX we extract the embedded image blob directly (no rendering).
        if describer is not None and len(text.split()) < IMAGE_WORD_THRESHOLD:
            cache_key = f"slide_{i}"
            if img_cache is not None and cache_key in img_cache:
                img_desc = img_cache[cache_key]
            else:
                img_desc = ""
                for shape in slide.shapes:
                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        try:
                            pil = PILImage.open(
                                io.BytesIO(shape.image.blob)).convert("RGB")
                            img_desc = describer.describe_slide_image(pil)
                            break   # one image per slide is enough
                        except Exception:
                            pass
                if img_cache is not None:
                    img_cache[cache_key] = img_desc

            if img_desc:
                text = (text + "\n" + img_desc).strip()

        slides.append(SlideText(i, label, text))
    return slides


def extract_docx(path: Path) -> list[SlideText]:
    """
    Word documents lack explicit page breaks; we treat each paragraph as one
    logical unit and group them into synthetic 'pages' of ≈ PAGE_PARA paragraphs.
    """
    from docx import Document
    PAGE_PARA = 15
    doc   = Document(str(path))
    paras = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    pages: list[SlideText] = []
    for pi, start in enumerate(range(0, max(len(paras), 1), PAGE_PARA)):
        chunk = paras[start:start + PAGE_PARA]
        text  = "\n".join(chunk)
        label = chunk[0][:80] if chunk else f"Page {pi+1}"
        pages.append(SlideText(pi, label, text))
    return pages


def load_slides(path: Path,
                describer: ImageDescriber | None = None,
                img_cache: dict | None = None) -> list[SlideText]:
    ext = path.suffix.lower()
    if ext == ".pdf":
        return extract_pdf(path, describer, img_cache)
    if ext in (".pptx", ".ppt"):
        return extract_pptx(path, describer, img_cache)
    if ext in (".docx", ".doc"):
        return extract_docx(path)
    raise ValueError(f"Unsupported slide format: {ext}")


# ── Embedding & FAISS index ───────────────────────────────────────────────────

def get_embedder():
    """Load sentence-transformer model once; use GPU if available."""
    from sentence_transformers import SentenceTransformer
    import torch
    device = "cuda" if torch.cuda.is_available() else "cpu"
    print(f"  [embed] Loading {EMBED_MODEL} on {device} ...")
    model = SentenceTransformer(EMBED_MODEL, device=device)
    return model


def embed_texts(model, texts: list[str]) -> np.ndarray:
    """Return L2-normalised float32 embeddings, shape (N, D)."""
    vecs = model.encode(
        texts,
        batch_size=BATCH_SIZE,
        show_progress_bar=False,
        normalize_embeddings=True,   # cosine sim → inner product on unit sphere
        convert_to_numpy=True,
    )
    return vecs.astype(np.float32)


def build_faiss_index(embeddings: np.ndarray) -> faiss.IndexFlatIP:
    """Inner-product index (= cosine similarity on normalised vectors)."""
    dim   = embeddings.shape[1]
    index = faiss.IndexFlatIP(dim)
    index.add(embeddings)
    return index


# ── Transcript windowing ──────────────────────────────────────────────────────

def build_window_texts(segments: list[dict], context_sec: float) -> list[str]:
    """
    For each segment build a richer query string by pooling the text of all
    segments whose midpoint falls within ±context_sec of this segment's midpoint.
    This prevents very short segments (e.g. "Yes" or "Okay") from misleading
    the matcher.
    """
    mids = [(s["start"] + s["end"]) / 2.0 for s in segments]
    n    = len(segments)
    texts: list[str] = []
    for i, mid in enumerate(mids):
        parts = []
        # scan backwards
        j = i
        while j >= 0 and mids[j] >= mid - context_sec:
            j -= 1
        # scan forwards
        k = i
        while k < n and mids[k] <= mid + context_sec:
            k += 1
        for s in segments[j+1 : k]:
            if s["text"].strip():
                parts.append(s["text"].strip())
        texts.append(" ".join(parts) if parts else (segments[i]["text"] or " "))
    return texts


# ── Viterbi temporal smoothing ────────────────────────────────────────────────

def viterbi_smooth(
    log_likelihoods: np.ndarray,   # shape (T, N_slides)
) -> list[int]:
    """
    Viterbi decoding that prefers forward progression.

    log_likelihoods[t, s]  = log P(observation_t | slide s)
                           = cosine_similarity score (already log-scale proxy)
    Returns: list of best slide indices, length T.
    """
    T, N = log_likelihoods.shape

    # dp[t, s] = best total log-score ending at slide s at step t
    dp   = np.full((T, N), -np.inf, dtype=np.float64)
    back = np.zeros((T, N), dtype=np.int32)

    dp[0] = log_likelihoods[0]

    for t in range(1, T):
        for s in range(N):
            # consider all previous states s_prev
            trans = np.zeros(N, dtype=np.float64)
            for sp in range(N):
                delta = s - sp
                if delta == 0:
                    trans[sp] = STAY_LOGP
                elif delta > 0:
                    trans[sp] = FWD_LOGP_PER * delta
                else:
                    trans[sp] = BWD_LOGP_PER * abs(delta)
            scores       = dp[t-1] + trans
            best_prev    = int(np.argmax(scores))
            dp[t, s]     = scores[best_prev] + log_likelihoods[t, s]
            back[t, s]   = best_prev

    # Traceback
    path     = [0] * T
    path[-1] = int(np.argmax(dp[-1]))
    for t in range(T - 2, -1, -1):
        path[t] = int(back[t + 1, path[t + 1]])
    return path


def viterbi_smooth_fast(log_likelihoods: np.ndarray) -> list[int]:
    """
    Vectorised Viterbi — O(T × N) instead of O(T × N²).
    Equivalent to the loop version but ~100× faster.
    """
    T, N = log_likelihoods.shape

    dp   = np.full((T, N), -np.inf, dtype=np.float64)
    back = np.zeros((T, N), dtype=np.int32)
    dp[0] = log_likelihoods[0]

    # Precompute transition matrix  trans[s_prev, s]
    idx   = np.arange(N)
    delta = idx[None, :] - idx[:, None]   # (N, N)  delta[sp, s] = s - sp
    trans = np.where(delta == 0, STAY_LOGP,
            np.where(delta >  0, FWD_LOGP_PER * delta,
                                 BWD_LOGP_PER * np.abs(delta)))

    for t in range(1, T):
        # scores[sp, s] = dp[t-1, sp] + trans[sp, s]
        scores        = dp[t-1, :, None] + trans          # (N, N)
        best_prev     = np.argmax(scores, axis=0)          # (N,)
        dp[t]         = scores[best_prev, np.arange(N)] + log_likelihoods[t]
        back[t]       = best_prev

    path     = [0] * T
    path[-1] = int(np.argmax(dp[-1]))
    for t in range(T - 2, -1, -1):
        path[t] = int(back[t + 1, path[t + 1]])
    return path


# ── Timeline collapse ─────────────────────────────────────────────────────────

def build_timeline(segments: list[dict], slide_path: list[int],
                   slides: list[SlideText],
                   off_slide_mask: list[bool] | None = None) -> list[dict]:
    """
    Merge consecutive segments assigned to the same slide into one interval.
    Off-slide segments are excluded from the timeline.
    Returns list of {slide_1based, start, end, label}.
    """
    if not segments:
        return []

    if off_slide_mask is None:
        off_slide_mask = [False] * len(segments)

    timeline: list[dict] = []
    cur_slide: int | None = None
    cur_start = 0.0
    cur_end   = 0.0

    for i in range(len(segments)):
        s  = segments[i]
        si = slide_path[i]

        if off_slide_mask[i]:
            # Flush current span before the gap
            if cur_slide is not None:
                timeline.append({
                    "slide":  cur_slide + 1,
                    "start":  round(cur_start, 3),
                    "end":    round(cur_end, 3),
                    "label":  slides[cur_slide].label,
                })
                cur_slide = None
            continue

        if cur_slide is None:
            cur_slide = si
            cur_start = s["start"]
            cur_end   = s["end"]
        elif si == cur_slide:
            cur_end = s["end"]
        else:
            timeline.append({
                "slide":  cur_slide + 1,
                "start":  round(cur_start, 3),
                "end":    round(cur_end, 3),
                "label":  slides[cur_slide].label,
            })
            cur_slide = si
            cur_start = s["start"]
            cur_end   = s["end"]

    if cur_slide is not None:
        timeline.append({
            "slide":  cur_slide + 1,
            "start":  round(cur_start, 3),
            "end":    round(cur_end, 3),
            "label":  slides[cur_slide].label,
        })
    return timeline


# ── Sparse-slide enrichment ───────────────────────────────────────────────────

def _enrich_sparse_slides(texts: list[str]) -> list[str]:
    """
    Slides that are nearly empty (section headers, diagram-only, code-only)
    get almost no embedding signal.  Enrich them by appending up to
    NEIGHBOR_WORDS from the nearest content-rich neighbours so the embedder
    has something to work with.  The original label text is kept at the front
    so the slide's own identity still dominates.
    """
    word_counts = [len(t.split()) for t in texts]
    enriched    = list(texts)

    for i, (text, wc) in enumerate(zip(texts, word_counts)):
        if wc >= SPARSE_THRESHOLD:
            continue
        extra: list[str] = []
        for delta in (1, -1, 2, -2, 3, -3):
            j = i + delta
            if 0 <= j < len(texts) and word_counts[j] >= SPARSE_THRESHOLD:
                words = texts[j].split()[:NEIGHBOR_WORDS]
                extra.extend(words)
            if len(extra) >= NEIGHBOR_WORDS * 2:
                break
        if extra:
            enriched[i] = text + " " + " ".join(extra)

    return enriched


# ── Main alignment routine ────────────────────────────────────────────────────

def align(caption_path: Path, slide_path: Path,
          out_dir: Path, embedder=None) -> Path:
    """
    Full alignment pipeline for one (caption, slide) pair.
    Returns path of the saved JSON.
    """
    print(f"\n{'='*70}")
    print(f"Caption : {caption_path.name}")
    print(f"Slides  : {slide_path.name}")
    print(f"{'='*70}")

    # ── Load inputs ──────────────────────────────────────────────────────────
    with open(caption_path, encoding="utf-8") as f:
        caption = json.load(f)
    segments: list[dict] = caption["segments"]
    if not segments:
        print("  [skip] Caption has no segments.")
        return out_dir

    print(f"  Transcript: {len(segments)} segments, {caption['duration']:.0f}s")

    # ── Load slides (with image captioning + speaker notes) ───────────────────
    describer = ImageDescriber()
    img_cache = _load_image_cache(slide_path)
    cache_size_before = len(img_cache)

    print("  Extracting slide text (+ images + speaker notes)...")
    slides = load_slides(slide_path, describer, img_cache)
    print(f"  Slides: {len(slides)} pages/slides")

    # Save image cache if any new descriptions were generated
    new_descriptions = len(img_cache) - cache_size_before
    if new_descriptions > 0:
        _save_image_cache(slide_path, img_cache)
        print(f"  [img] {describer.api_calls} GPT-4o-mini vision API call(s) made "
              f"({new_descriptions} new slide(s) described, cached for future runs)")

    # ── Embed slides ──────────────────────────────────────────────────────────
    if embedder is None:
        embedder = get_embedder()

    slide_texts = _enrich_sparse_slides([s.text if s.text.strip() else s.label for s in slides])
    print(f"  Embedding {len(slides)} slides...")
    slide_embs = embed_texts(embedder, slide_texts)
    index      = build_faiss_index(slide_embs)

    # ── Embed transcript segments (with context window) ───────────────────────
    print(f"  Building context windows (±{CONTEXT_SEC:.0f}s)...")
    window_texts = build_window_texts(segments, CONTEXT_SEC)

    print(f"  Embedding {len(window_texts)} transcript windows...")
    seg_embs = embed_texts(embedder, window_texts)   # (T, D)

    # ── Raw similarity scores: each segment vs all slides ─────────────────────
    # faiss.search returns (scores, indices); with IndexFlatIP + normalised
    # vectors scores are cosine similarities in [-1, 1].
    print("  Querying FAISS index...")
    k          = len(slides)
    sims, idxs = index.search(seg_embs, k)   # both (T, N_slides)

    # Build log-likelihood matrix: reorder sims by slide index
    T          = len(segments)
    N          = len(slides)
    log_ll     = np.zeros((T, N), dtype=np.float64)
    for t in range(T):
        for rank in range(k):
            log_ll[t, idxs[t, rank]] = float(sims[t, rank])

    # Save raw log-likelihoods for off-slide detection (before the prior)
    log_ll_raw = log_ll.copy()

    # ── Temporal position prior ───────────────────────────────────────────────
    # At segment t, the professor is expected to be near slide
    # expected_s = (t_mid / total_duration) * (N - 1).
    # A Gaussian prior around this expected position gently pushes the
    # Viterbi toward the correct slide when raw similarities are ambiguous
    # (e.g., consecutive slides with nearly identical content).
    total_duration = caption["duration"] or 1.0
    slide_idx      = np.arange(N, dtype=np.float64)
    for t, seg in enumerate(segments):
        t_mid      = (seg["start"] + seg["end"]) / 2.0
        expected_s = (t_mid / total_duration) * (N - 1)
        prior      = -0.5 * ((slide_idx - expected_s) / PRIOR_SIGMA) ** 2
        log_ll[t] += prior

    # ── Off-slide detection ───────────────────────────────────────────────────
    # Segments where the best raw cosine is below threshold are flagged.
    # The professor may be speaking off-topic, answering questions, or doing
    # live demos that don't correspond to any slide.
    raw_max_sim    = log_ll_raw.max(axis=1)   # (T,)
    off_slide_mask = (raw_max_sim < OFF_SLIDE_THRESHOLD).tolist()
    n_off = sum(off_slide_mask)
    if n_off:
        print(f"  Off-slide segments: {n_off}/{T} "
              f"(raw cosine < {OFF_SLIDE_THRESHOLD})")

    # ── Viterbi ───────────────────────────────────────────────────────────────
    print("  Running Viterbi smoothing...")
    slide_path_idx = viterbi_smooth_fast(log_ll)

    # ── Per-segment results ───────────────────────────────────────────────────
    aligned_segments = []
    for i, seg in enumerate(segments):
        si        = slide_path_idx[i]
        off       = off_slide_mask[i]
        raw_sim   = round(float(log_ll_raw[i, si]), 4)
        aligned_segments.append({
            "id":          seg["id"],
            "start":       seg["start"],
            "end":         seg["end"],
            "text":        seg["text"],
            "slide":       None if off else si + 1,    # 1-based; null = off-slide
            "slide_label": None if off else slides[si].label,
            "similarity":  raw_sim,
            "off_slide":   off,
        })

    timeline = build_timeline(segments, slide_path_idx, slides, off_slide_mask)

    # ── Output ────────────────────────────────────────────────────────────────
    result = {
        "lecture":        caption_path.stem,
        "slide_file":     slide_path.name,
        "total_slides":   len(slides),
        "total_segments": len(segments),
        "off_slide_count": n_off,
        "duration":       caption["duration"],
        "language":       caption.get("language", ""),
        "embed_model":    EMBED_MODEL,
        "context_sec":    CONTEXT_SEC,
        "segments":       aligned_segments,
        "timeline":       timeline,
    }

    out_dir.mkdir(parents=True, exist_ok=True)
    out_file = out_dir / f"{caption_path.stem}.json"
    with open(out_file, "w", encoding="utf-8") as f:
        json.dump(result, f, ensure_ascii=False, indent=2)

    # Summary
    slide_counts = {}
    for i, si in enumerate(slide_path_idx):
        if not off_slide_mask[i]:
            slide_counts[si] = slide_counts.get(si, 0) + 1
    covered = len(slide_counts)
    print(f"  Covered {covered}/{N} slides across {len(timeline)} intervals")
    print(f"  Saved → {out_file}")
    return out_file


# ── Auto-discovery ────────────────────────────────────────────────────────────

def _candidate_slides(course_dir: Path) -> list[Path]:
    """All PDF/PPTX/DOCX under [course]/materials/."""
    exts = {".pdf", ".pptx", ".ppt", ".docx", ".doc"}
    return [p for p in (course_dir / "materials").rglob("*") if p.suffix.lower() in exts]


def _name_similarity(a: str, b: str) -> float:
    """Rough token-overlap score between two filenames (lowercase, no ext)."""
    ta = set(a.lower().replace("-", " ").replace("_", " ").split())
    tb = set(b.lower().replace("-", " ").replace("_", " ").split())
    if not ta or not tb:
        return 0.0
    return len(ta & tb) / len(ta | tb)


def find_best_slide(caption_path: Path, slide_candidates: list[Path]) -> Path | None:
    """Return the slide file most likely corresponding to the caption."""
    stem = caption_path.stem
    best_score, best = 0.0, None
    for sp in slide_candidates:
        score = _name_similarity(stem, sp.stem)
        if score > best_score:
            best_score, best = score, sp
    return best if best_score > 0.05 else None


def process_course(course_id: int | str) -> None:
    course_dir = PROJECT_DIR / str(course_id)
    captions   = list((course_dir / "captions").glob("*.json"))
    slides     = _candidate_slides(course_dir)
    out_dir    = course_dir / "alignment"

    if not captions:
        print(f"No captions found in {course_dir}/captions/")
        return
    if not slides:
        print(f"No slide files found under {course_dir}/materials/")
        return

    embedder = get_embedder()

    for cap in captions:
        out_file = out_dir / f"{cap.stem}.json"
        if out_file.exists():
            print(f"  [skip] Already aligned: {cap.stem}")
            continue

        slide = find_best_slide(cap, slides)
        if slide is None:
            print(f"  [warn] No matching slide for {cap.name} — skipping")
            continue

        align(cap, slide, out_dir, embedder=embedder)


# ── Entry point ───────────────────────────────────────────────────────────────

def main() -> None:
    parser = argparse.ArgumentParser(description="Semantic caption↔slide alignment")
    parser.add_argument("--caption", metavar="PATH", help="Whisper caption JSON")
    parser.add_argument("--slides",  metavar="PATH", help="Slide file (PDF/PPTX/DOCX)")
    parser.add_argument("--course",  metavar="ID",   help="Auto-process a course folder")
    parser.add_argument("--out",     metavar="DIR",  help="Output directory (default: [course]/alignment)")
    args = parser.parse_args()

    if args.course:
        process_course(args.course)
        return

    if not args.caption or not args.slides:
        parser.error("Provide both --caption and --slides, or use --course ID")

    cap_path   = Path(args.caption)
    slide_path = Path(args.slides)

    if not cap_path.exists():
        print(f"[error] Caption not found: {cap_path}"); sys.exit(1)
    if not slide_path.exists():
        print(f"[error] Slides not found: {slide_path}"); sys.exit(1)

    # Infer output dir from caption path: two levels up → alignment/
    if args.out:
        out_dir = Path(args.out)
    else:
        out_dir = cap_path.parent.parent / "alignment"

    embedder = get_embedder()
    align(cap_path, slide_path, out_dir, embedder=embedder)


if __name__ == "__main__":
    main()
