"""
Note Generation
Generates one comprehensive Markdown note file per course, covering all
lectures in sequence, matching the style of example/CS2105_note.md.

Architecture:
  - Per lecture: split slides into ~CHAPTER_SIZE chunks, one GPT call per chunk.
  - Chunks map to ### N.x sections. Images injected at diagram slides.
  - All lectures merged into one file; exam notes appended at the end.
  - Self-scoring via heuristics (no extra API call).

Usage:
  python note_generation.py --course 85427
  python note_generation.py --course 85427 --detail 9 --iterate
  python note_generation.py --slides 85427/materials/LectureNotes/L02.pdf \\
      --alignment "85427/alignment/L02.json" --lecture-num 2 --course-name "CS3210"
"""

from __future__ import annotations

import argparse
import json
import re
import sys
from datetime import datetime
from pathlib import Path

try:
    from tqdm import tqdm
except ImportError as _e:
    print(f"[error] Missing dependency: {_e}")
    print("[error] Please install the ML environment from Settings → ML Environment in the AutoNote app.")
    sys.exit(1)

try:
    import alignment_parser
except ImportError as _e:
    print(f"[error] Could not import alignment_parser: {_e}")
    print(f"[error] Make sure alignment_parser.py is in the same directory as note_generation.py")
    sys.exit(1)

PROJECT_DIR = Path(__file__).parent
import sys as _sys
_AUTO_NOTE_DIR = Path.home() / ".auto_note"
import os as _os
if _os.environ.get("AUTONOTE_DATA_DIR"):
    DATA_DIR = Path(_os.environ["AUTONOTE_DATA_DIR"])
elif getattr(_sys, "frozen", False) or PROJECT_DIR == _AUTO_NOTE_DIR / "scripts":
    DATA_DIR = _AUTO_NOTE_DIR
else:
    DATA_DIR = PROJECT_DIR

# Course output directory: defaults to DATA_DIR but can be overridden by
# OUTPUT_DIR in config.json so files land in the user's chosen Output Dir.
_ng_config: dict = (
    json.loads((DATA_DIR / "config.json").read_text())
    if (DATA_DIR / "config.json").exists() else {}
)
_out_dir = _ng_config.get("OUTPUT_DIR", "").strip()
COURSE_DATA_DIR = Path(_out_dir) if _out_dir else Path.home() / "AutoNote"

# ── Constants ─────────────────────────────────────────────────────────────────

DETAIL_LEVEL      = 7
OUTPUT_FORMAT     = "md"
NOTE_MODEL        = "gpt-5.1"
TRANSLATE_MODEL   = "gpt-4o"   # Chinese/other translation post-pass — gpt-4o is cheap enough
QUALITY_TARGET    = 8.0
IMAGE_RENDER_SCALE = 1.5
NOTE_LANGUAGE     = "en"    # "en" = English | "zh" = Chinese
SHOW_SCORE        = False   # dev mode: set via --score flag to show self-scoring

CHAPTER_SIZE      = 15      # slides per GPT call
MAX_NOTE_CHARS = 120000  # max total chars in the chunk prompt (transcript + slides + images)

SCORE_WEIGHTS = {"coverage": 0.30, "terminology": 0.35,
                 "callouts": 0.15, "code_blocks": 0.20}
MIN_NOTE_WORDS_PER_SLIDE = 60   # expected words per slide in final note


# ── Prompts ───────────────────────────────────────────────────────────────────

_PROMPTS: dict[str, dict] = {

"en": dict(
system="""\
You are a teaching assistant at a top university, writing high-quality study notes for computer science courses based on lecture slides and audio transcripts.

Writing guidelines:
1. Write in English. Keep technical terms in English.
2. Never use a third-person narrator perspective. Do not write "the professor said", "the lecturer pointed out", etc. Focus on the knowledge itself — state concepts, principles, and conclusions directly:
   - ✗ "The professor explained that…"  →  ✓ "The key idea is…"
   - ✗ "The lecturer used an example…"  →  ✓ "As an example,…"
3. Structure content as: concept → principle → example → exam focus. Write fluent explanatory paragraphs; do not list slide bullets verbatim.
4. Use LaTeX for math: inline $...$, display $$...$$.
5. Code examples must be complete, compilable/runnable snippets (with necessary includes, function signatures, main, etc.) using correct syntax highlighting (```c, ```cpp, ```python, etc.). Use pseudocode only when no real equivalent exists, tagged as ```pseudo.
6. Mark exam-critical content with:
     > [!IMPORTANT]
     > content
7. Use italics for interesting analogies or memory aids.
8. Image insertion rules (strictly follow):
   - **Insert all and only images that contain visual elements**: diagrams, flowcharts, architecture drawings, code screenshots, mathematical derivations, data visualizations, tables with meaningful structure, annotated figures, or any non-trivial visual illustration. Do NOT insert administrative or non-course elements (course info slides, polling QR codes, attendance prompts, etc.) even if they contain images.
   - Pure text slides (bullet points, definitions, titles) do not need images — the notes express text better than a screenshot.
   - **Be INCLUSIVE with images**: if a frame/slide shows a diagram, chart, table, code, or any non-trivial visual content, include it. Aim to include most of the content-rich images available, not just a few highlights. It's better to have more images with brief connecting text than to skip images.
   - **Each image MUST be placed inline, immediately after the paragraph that directly discusses the concept shown in that image.** If a frame shows content that the transcript doesn't fully cover, briefly describe the frame's content (using the description provided in the image hints) and then insert the image.
   - Format for slide images: `![Slide N](images/L**…/slide_NNN.png) *(one-sentence description)*`
   - Format for screen-capture frames: `![Frame N](images/L**…/frame_NNN.png) *(one-sentence description)*`
     (The subdirectory under `images/` is provided in the "Available images" list — copy it verbatim, including any slug suffix, and do not shorten or rewrite it. The caption must be in parentheses wrapped in asterisks exactly as shown.)
9. Never fabricate technical details not present in the source material.
""",
chunk="""\
Write study notes for the following course segment ({course_name} Lecture {lec_num}: {lec_title}).

## Lecture audio transcript (PRIMARY SOURCE — this is the main content to cover)
{transcript_block}

## Slide outline (structural guide — use for topic organization)
{slide_outline}

## Available images (insert relevant ones inline)
{image_hints}

---

Requirements:
- The **transcript is the primary source material**. Cover ALL concepts, explanations, examples, and details the lecturer discusses. The slide outline is a structural guide for organizing topics, but the transcript contains the actual teaching content.
- The section heading for this segment is `### {lec_num}.{chunk_idx} {chunk_title}` (**do not output this line** — it is added by the caller).
- Detail level: {detail}/10. {detail_instruction}
- Images: **insert all images that contain visual elements** (diagrams, charts, graphs, code screenshots, architecture drawings, data visualizations, mathematical derivations, etc.). Skip images of pure text, bullet points, or administrative/non-course elements.
  Copy the exact path from the "Available images" list above (including the images/L** subdirectory). Do not invent paths.
  **CRITICAL: Be inclusive — aim to include MOST content-rich images (diagrams, charts, code, tables). For each image, ensure a paragraph discusses its content, then insert the image right after that paragraph. If a frame has visual content that the transcript doesn't cover, write a brief paragraph about it based on the image description, then insert the image. NEVER cluster multiple images together consecutively without explanatory text between them.**
  Format: `![Slide N](path) *(caption)*` or `![Frame N](path) *(caption)*`
- Code examples must be complete and compilable (with necessary includes/imports), using the correct language tag (```c, ```cpp, ```python).
- Only cover the content in this segment; do not introduce material from other lectures.
""",
slide_only="""\
Write study notes for {course_name} Lecture {lec_num}: {lec_title} based on the slides below.
(No audio transcript is available — supplement with your CS knowledge where appropriate.)

## Slide content
{slide_outline}

## Available images
{image_hints}

---

Requirements:
- The section heading is `### {lec_num}.{chunk_idx} {chunk_title}` (**do not output this line**).
- Detail level: {detail}/10. {detail_instruction}
- Images: **insert all images that contain visual elements** (diagrams, charts, graphs, code screenshots, architecture drawings, data visualizations, mathematical derivations, etc.). Skip images of pure text, bullet points, or administrative/non-course elements.
  Copy the exact path from the "Available images" list above (including the images/L** subdirectory). Do not invent paths.
  **CRITICAL: Be inclusive — aim to include MOST content-rich images. For each image, write a paragraph about its content then insert the image. NEVER cluster multiple images together consecutively without explanatory text between them.**
  Format: `![Slide N](path) *(caption)*` or `![Frame N](path) *(caption)*`
- Code examples must be complete and compilable, using the correct language tag (```c, ```cpp, ```python).
""",
exam="""\
Below are the complete lecture notes for {course_name}. Please append a concise exam cheat-sheet section at the end.

Format:
- Heading: `## Exam Notes`
- Each entry: `N. **Topic**: one-sentence summary`
- No more than 30 entries, covering key concepts, formulas, algorithm steps, and common confusion points from all lectures.

Notes summary:
{summary}
""",
no_transcript="(No audio transcript available for this segment.)",
detail_instructions=[
    (range(0, 3),  "Minimal bullets: one line per concept, no expansion, max 3 bullets per slide."),
    (range(3, 6),  "Hierarchical bullets: one top-level bullet (`-`) per main concept, "
                   "at most 2 sub-bullets (`  -`) for key details. "
                   "Max 5 bullets total per slide. No prose paragraphs."),
    (range(6, 9),  "Detailed paragraphs: cover concepts, principles, the lecturer's examples and analogies in full."),
    (range(9, 11), "Maximum detail: include all nuances, edge cases, connections to other chapters, and exam pointers."),
],
),  # end en

}  # end _PROMPTS


# Language names for the translation instruction
_LANG_NAMES = {"en": "English", "zh": "Chinese", "ja": "Japanese", "ko": "Korean"}


def _P(key: str) -> str:
    """Return the English prompt unchanged. Translation is always a separate
    post-generation step via _translate()."""
    return _PROMPTS["en"][key]


def _detail_instr(level: int) -> str:
    instrs = _PROMPTS["en"]["detail_instructions"]
    for rng, txt in instrs:
        if level in rng:
            return txt
    return instrs[2][1]

# ── Image filter constants ────────────────────────────────────────────────────

IMAGE_FILTER_MODEL      = "gpt-4o"
IMAGE_FILTER_WORD_MAX   = 12   # slides with ≤ this many words → remove without API call
IMAGE_FILTER_HEURISTIC  = 80   # slides with > this many words AND no code/desc → remove

# Title/divider patterns that add no visual value
_TITLE_PATTERN = re.compile(
    r"^\s*(CS\d+|AY\d+|Lecture\s+\d+|\[.*\]|Part\s+\d+|Section\s+\d+|"
    r"Outline|Agenda|Table of Contents|Overview|Summary|Questions\?|Q&A)\s*$",
    re.IGNORECASE | re.MULTILINE,
)

# Keywords that indicate a slide description contains visual/diagram content
_VISUAL_KEYWORDS = re.compile(
    r"\b(diagram|chart|graph|figure|illustration|flowchart|screenshot|"
    r"table|formula|equation|architecture|layout|structure|matrix|tree|"
    r"network|circuit|timeline|image|photo|plot|drawing|schematic|visual)\b",
    re.IGNORECASE,
)


def _desc_has_visual(desc: str) -> bool:
    """Return True if a cache description mentions diagram/chart/visual content."""
    return bool(_VISUAL_KEYWORDS.search(desc))


def _img_ref_pattern() -> re.Pattern:
    # Matches slide_001.png / frame_001.png under images/L04/, images/L04_F02/,
    # or images/L04_<slug>[_F02]/ — with an optional trailing italic caption.
    return re.compile(
        r"!\[(?:Slide|Frame) \d+\]\((images/L\d{2}[^/]*/(?:slide|frame)_\d{3}\.png)\)"
        r"(?:\s*\*\([^)]*\)\*)?"
    )


def _vision_keep(img_path: Path, slide_text: str = "") -> bool:
    """Ask GPT-4o-mini whether the slide image is worth including in notes.

    Uses slide_text as context so the model can reason about whether the
    visual structure actually explains the lecture content (e.g. a diagram
    of a protocol stack/a picture of a biological structure) vs. being an unrelated administrative element
    (e.g. a PollEv QR code, a course logo, a participation prompt).
    """
    import base64
    import io
    if not img_path.exists():
        return False
    try:
        from PIL import Image as PILImage
        img = PILImage.open(img_path).convert("RGB")
        # Downscale to max 800px wide to keep base64 payload small
        if img.width > 800:
            ratio = 800 / img.width
            img = img.resize((800, int(img.height * ratio)), PILImage.LANCZOS)
        buf = io.BytesIO()
        img.save(buf, format="JPEG", quality=75)
        b64 = base64.b64encode(buf.getvalue()).decode("ascii")
    except Exception:
        return True   # can't load image → default keep

    context_block = (
        f"\n\nSlide text (OCR):\n\"\"\"\n{slide_text[:400]}\n\"\"\""
        if slide_text.strip() else ""
    )

    _VISION_PROMPT = f"""\
You are a study-notes curator deciding whether a lecture slide image should be \
embedded in written notes.{context_block}

## KEEP the image if the slide contains ANY course-relevant visual element:
   - Diagrams: system/architecture diagrams, component boxes connected by arrows
   - Flowcharts, state machines, decision trees, sequence/timing diagrams
   - Memory layouts, address-space maps, cache/pipeline stage illustrations
   - Graphs, plots, bar/line/pie charts, scatter plots showing data or trends
   - Tables with a meaningful grid structure (comparing options, relationships)
   - Mathematical formulas or derivations where spatial layout matters
   - Code screenshots or annotated code with visual highlights
   - Annotated screenshots, highlighted output, or callout arrows
   - Any non-trivial visual illustration related to the course

## REMOVE only if the slide is clearly non-visual or non-course-related:
- Pure text slides (bullet points, prose, definitions) with absolutely no \
diagram, chart, figure, or visual element
- Title slides, section dividers, agenda/outline, blank slides
- Administrative elements unrelated to the course: polling/quiz prompts \
(PollEv, Mentimeter, Kahoot QR codes), attendance check slides, \
course info graphics, "any questions?" slides, logos, sponsor slides

## Default: when genuinely uncertain, KEEP.

Reply with exactly one word: KEEP or REMOVE."""

    try:
        openai_client = _get_client_for(IMAGE_FILTER_MODEL)
        r = openai_client.chat.completions.create(
            model=IMAGE_FILTER_MODEL,
            messages=[{"role": "user", "content": [
                {"type": "image_url",
                 "image_url": {"url": f"data:image/jpeg;base64,{b64}", "detail": "low"}},
                {"type": "text", "text": _VISION_PROMPT},
            ]}],
            max_tokens=5,
        )
        return "KEEP" in r.choices[0].message.content.strip().upper()
    except Exception:
        return True   # default: keep on API error


def filter_images_pass(
    notes_text: str,
    notes_dir: Path,
    lectures: list["LectureData"],
) -> tuple[str, int, int]:
    """Post-processing agent: remove low-value image references from merged notes.

    Decision priority:
      1. Cache-verified AND description mentions visual elements → KEEP
      2. Title/divider pattern → REMOVE
      3. All other cases → vision API decision

    Returns (cleaned_text, n_kept, n_removed).
    """
    pattern = _img_ref_pattern()

    # Build unified lookup: image rel-path → (SlideInfo, LectureData).
    # Path scheme matches render_chunk_images: images/{dir_key}/…, where
    # dir_key = L{num:02d}_{slug}[_F{idx:02d}].
    slide_ld_lookup: dict[str, tuple[SlideInfo, "LectureData"]] = {}
    for ld in lectures:
        # Use frame_NNN for screenshare, slide_NNN for traditional
        img_prefix = "frame" if ld.source == "screenshare" else "slide"
        for s in ld.slides:
            key = f"images/{ld.dir_key}/{img_prefix}_{s.index+1:03d}.png"
            slide_ld_lookup[key] = (s, ld)

    # Collect unique paths and decide keep/remove
    decisions: dict[str, bool] = {}   # path → True=keep
    for m in pattern.finditer(notes_text):
        rel = m.group(1)
        if rel in decisions:
            continue

        pair    = slide_ld_lookup.get(rel)
        slide   = pair[0] if pair else None
        owner   = pair[1] if pair else None
        img_path = notes_dir / rel

        # ⓪ Screen share frames — always keep (the frame IS the content)
        if owner and owner.source == "screenshare":
            decisions[rel] = True
            continue

        # ① Cache-verified AND description mentions visual elements → KEEP
        if slide and owner:
            desc = owner.img_cache.get(f"page_{slide.index}", "")
            if desc and _desc_has_visual(desc):
                decisions[rel] = True
                continue

        # ② Title/divider pattern → REMOVE
        if slide and _TITLE_PATTERN.search(slide.text):
            decisions[rel] = False
            continue

        # ③ Vision API — KEEP only if visual AND relevant to lecture content
        slide_text = slide.text if slide else ""
        decisions[rel] = _vision_keep(img_path, slide_text)

    kept    = sum(1 for v in decisions.values() if v)
    removed = sum(1 for v in decisions.values() if not v)
    tqdm.write(f"  Image filter: {kept} kept, {removed} removed out of {len(decisions)}")

    # Remove lines for filtered-out images; collapse extra blank lines
    lines_out: list[str] = []
    for line in notes_text.splitlines():
        m = pattern.fullmatch(line.strip())
        if m and not decisions.get(m.group(1), True):
            # Replace filtered image line with nothing (don't emit the line)
            continue
        lines_out.append(line)

    cleaned = "\n".join(lines_out)
    # Collapse 3+ consecutive blank lines → 2 (preserves paragraph spacing)
    cleaned = re.sub(r"\n{3,}", "\n\n", cleaned)
    return cleaned, kept, removed


def _max_tokens(level: int) -> int:
    # Per chunk (CHAPTER_SIZE slides).
    # gpt-5.x reasoning models consume tokens for internal thinking,
    # so we need larger budgets than the expected output length.
    # Token budget directly caps output length — keep it proportional to detail level.
    if level < 3:  return 2000
    if level < 6:  return 3500
    if level < 9:  return 10000
    return 16000


# ── Multi-provider LLM helpers ────────────────────────────────────────────────

def _provider(model: str) -> str:
    if model == "claude-cli":
        return "claude-cli"
    if model == "codex-cli":
        return "codex-cli"
    if model.startswith("gemini"):
        return "gemini"
    if model.startswith("claude"):
        return "anthropic"
    if model.startswith("deepseek"):
        return "deepseek"
    if model.startswith("grok"):
        return "grok"
    if model.startswith(("mistral", "codestral", "pixtral", "magistral")):
        return "mistral"
    return "openai"


_client_cache: dict = {}


def _make_client(provider: str):
    import os
    from openai import OpenAI

    def _read_key(env_var: str, filename: str, label: str) -> str:
        key = os.environ.get(env_var, "")
        if not key:
            kf = DATA_DIR / filename
            if kf.exists():
                key = kf.read_text().strip()
        if not key:
            raise RuntimeError(
                f"No {label} API key found "
                f"(set {filename} in ~/.auto_note/ or {env_var} env var)"
            )
        return key

    if provider == "gemini":
        return OpenAI(
            api_key=_read_key("GEMINI_API_KEY", "gemini_api.txt", "Gemini"),
            base_url="https://generativelanguage.googleapis.com/v1beta/openai/",
        )
    elif provider == "anthropic":
        from anthropic import Anthropic
        return Anthropic(api_key=_read_key("ANTHROPIC_API_KEY", "anthropic_key.txt", "Anthropic"))
    elif provider == "deepseek":
        return OpenAI(
            api_key=_read_key("DEEPSEEK_API_KEY", "deepseek_key.txt", "DeepSeek"),
            base_url="https://api.deepseek.com",
        )
    elif provider == "grok":
        return OpenAI(
            api_key=_read_key("GROK_API_KEY", "grok_key.txt", "xAI Grok"),
            base_url="https://api.x.ai/v1",
        )
    elif provider == "mistral":
        return OpenAI(
            api_key=_read_key("MISTRAL_API_KEY", "mistral_key.txt", "Mistral"),
            base_url="https://api.mistral.ai/v1",
        )
    else:  # openai
        return OpenAI(api_key=_read_key("OPENAI_API_KEY", "openai_api.txt", "OpenAI"))


def _get_client_for(model: str):
    """Return (and cache) the appropriate API client for the given model name."""
    p = _provider(model)
    if p not in _client_cache:
        _client_cache[p] = _make_client(p)
    return _client_cache[p]


def _pick_translate_model() -> str:
    """Choose a translator that won't truncate. When the user's NOTE_MODEL
    has a much bigger output cap than TRANSLATE_MODEL (e.g. deepseek-v4-*
    has 384K, gpt-5.1 has 128K, gpt-4o only 16K), prefer NOTE_MODEL so
    long Chinese translations of long English drafts don't get cut off
    mid-sentence."""
    if NOTE_MODEL in ("claude-cli", "codex-cli"):
        return NOTE_MODEL
    note_cap = _MODEL_MAX_COMPLETION.get(NOTE_MODEL, 0)
    tr_cap = _MODEL_MAX_COMPLETION.get(TRANSLATE_MODEL, 0)
    if note_cap and tr_cap and note_cap > tr_cap * 2:
        return NOTE_MODEL
    return TRANSLATE_MODEL


def _translate(text: str, lang: str) -> str:
    """Translate note text to the target language, preserving all Markdown
    formatting, image references, LaTeX formulas, and code blocks verbatim.
    Only prose text is translated; technical terms keep English with
    translation in parentheses on first use.

    On detected truncation (finish_reason=length), the text is split at
    paragraph boundaries and translated chunk-by-chunk; the chunks are
    concatenated back together. This avoids the silent-truncation bug
    where a long English draft turned into a short Chinese fragment that
    ended mid-sentence or mid-image-link."""
    system = (
        f"You are a professional translator for technical study notes. "
        f"Translate English prose into {lang} while keeping ALL technical "
        f"terminology in its original English form. Do NOT translate "
        f"technical terms — readers are studying the subject in English "
        f"and need to recognize the exact English terminology from lectures, "
        f"exams, and textbooks."
    )
    prompt = (
        f"Translate the following study notes into {lang}.\n\n"
        f"Rules:\n"
        f"1. Translate ONLY the connecting prose (explanatory sentences, "
        f"narrative text) into {lang}.\n"
        f"2. Keep ALL technical terminology in ENGLISH, verbatim. This "
        f"includes — but is not limited to — protocol names (TCP, UDP, HTTP, "
        f"DHCP, ARP, ICMP, DNS, RSA, AES…), networking concepts (subnet mask, "
        f"MAC address, broadcast, unicast, frame, packet, segment, hub, "
        f"switch, router, bridge, LAN, WAN, VPN…), cryptography terms "
        f"(symmetric key, public key, private key, cipher, plaintext, "
        f"ciphertext, block cipher, stream cipher, hash, signature, session "
        f"key, Diffie-Hellman…), algorithm names (CSMA/CD, Caesar cipher, "
        f"monoalphabetic cipher, polyalphabetic cipher), proper nouns "
        f"(Alice, Bob, Trudy, Ethernet, Wi-Fi, OSI, NUS…), and anything in "
        f"code font `…`. Never translate these into {lang}.\n"
        f"3. Do NOT write the {lang} translation next to the English term — "
        f"just keep the English term as-is. The reader already understands "
        f"English technical vocabulary; they need {lang} only for the "
        f"connecting narrative.\n"
        f"4. Keep EXACTLY as-is without any modification:\n"
        f"   - Image lines: ![Slide N](path) *(caption)* — translate ONLY "
        f"the non-technical prose inside *(...)*, keep the path + all "
        f"English technical terms unchanged\n"
        f"   - LaTeX: $...$ and $$...$$\n"
        f"   - Code blocks: ```...```\n"
        f"   - Callout markers: > [!IMPORTANT]\n"
        f"   - Markdown formatting: ###, **, *, ---, etc.\n"
        f"5. Do NOT shorten, summarize, or omit any content.\n"
        f"6. Output ONLY the translated text.\n\n"
        f"Example (for Chinese):\n"
        f"  IN:  The symmetric key cryptography scheme uses the same key "
        f"for encryption and decryption.\n"
        f"  OUT: symmetric key cryptography 方案在 encryption 和 decryption "
        f"时使用同一个 key。\n\n"
        f"---\n\n{text}"
    )
    _tmodel = _pick_translate_model()
    # Generous output budget — Chinese translation of English text often
    # tokenizes 1.5-2x larger than the source on cl100k_base, so naive
    # len(text)*3 still bumps gpt-4o's 16K cap on chunks past ~5K chars.
    _trunc: list[bool] = []
    out = _call(_tmodel, system, prompt, len(text) * 3, _truncated=_trunc)
    flagged = _trunc and _trunc[0]
    # Some providers (notably DeepSeek's V4 chat completions) return
    # finish_reason="stop" even when the output was cut mid-token. Fall
    # back to a content-shape heuristic — broken UTF-8, mid-word cuts,
    # unbalanced markdown image links — so those silent truncations
    # don't end up in the cache.
    if not flagged and _looks_truncated(out):
        flagged = True
    if flagged:
        # Split at paragraph boundaries and translate chunk-by-chunk to
        # stay under the per-call output cap. Falls back to keeping the
        # English source for any chunk that still won't fit, rather than
        # caching a truncated translation.
        return _translate_chunked(text, lang)
    return out


_TRUNC_IMAGE_RE = re.compile(r"!\[[^\]]*\]\([^)\s\n]*\Z")
_SENTENCE_END_RE = re.compile(
    r"[.!?。！？]\s*[*_`>]*\s*\Z|[\)\]\*_`>]\s*\Z|[一-鿿][\)\]\*_`>]?\s*\Z"
    # Permissive: end with western/CJK terminator, closing markdown
    # punctuation, or any Chinese character (Chinese sentences often
    # end with the period folded into the last char's metric).
)


def _looks_truncated(text: str) -> bool:
    """Heuristic truncation detector for cases where the provider returns
    finish_reason='stop' but the text was actually cut mid-stream.

    Conservative — false positives waste a chunk-translate retry but
    don't lose data; false negatives cache a broken section, which is
    what we just shipped a fix for.
    """
    if not text:
        return True
    s = text.rstrip()
    if not s:
        return True
    # 1. Broken UTF-8 replacement char at the end → certain truncation
    if s.endswith("�"):
        return True
    # 2. Open image-link without closing paren — `![alt](path` cut
    if _TRUNC_IMAGE_RE.search(s):
        return True
    # 3. Ends mid-ASCII-word (alphabetic char, no sentence end nearby)
    last = s[-1]
    if last.isascii() and last.isalpha():
        # Allow technical term endings ONLY if a closing punct sits
        # within the last few chars. `_SENTENCE_END_RE` matches the
        # tail with permissive markdown closings; if it doesn't match
        # AND we end on a bare letter, assume mid-word truncation.
        if not _SENTENCE_END_RE.search(s[-12:]):
            return True
    return False


def _translate_chunked(text: str, lang: str, max_chunk_chars: int = 3500) -> str:
    """Recursive paragraph-by-paragraph translation. Joined back with the
    same separator the splitter used so Markdown structure is preserved."""
    paragraphs = text.split("\n\n")
    out_parts: list[str] = []
    cur: list[str] = []
    cur_len = 0
    for p in paragraphs:
        if cur_len + len(p) + 2 > max_chunk_chars and cur:
            out_parts.append("\n\n".join(cur))
            cur = [p]
            cur_len = len(p)
        else:
            cur.append(p)
            cur_len += len(p) + 2
    if cur:
        out_parts.append("\n\n".join(cur))

    translated: list[str] = []
    for part in out_parts:
        if not part.strip():
            translated.append(part)
            continue
        try:
            t = _translate(part, lang)   # depth-limited: chunks are small
            translated.append(t)
        except Exception:
            # Failed to translate this chunk — keep it in English rather
            # than dropping content silently.
            translated.append(part)
    return "\n\n".join(translated)


_MODEL_MAX_COMPLETION = {
    # Conservative per-model output-token caps for OpenAI models. The API
    # rejects requests where max_tokens exceeds these, so we clamp here.
    "gpt-4o":          16384,
    "gpt-4o-2024-08-06": 16384,
    "gpt-4o-2024-11-20": 16384,
    "gpt-4o-mini":     16384,
    "gpt-4.1":         32768,
    "gpt-4.1-mini":    32768,
    "gpt-4.1-nano":    32768,
    "gpt-5.1":         128000,
    "gpt-5.2":         128000,
    "o3":              100000,
    "o4-mini":         100000,
    # DeepSeek V4 (Pro + Flash) share a 384K max-output cap on the public
    # API; legacy deepseek-chat / deepseek-reasoner alias to v4-flash.
    "deepseek-v4-pro":   384000,
    "deepseek-v4-flash": 384000,
    "deepseek-chat":     384000,
    "deepseek-reasoner": 384000,
}


def _cap_tokens(model: str, max_tokens: int) -> int:
    """Clamp max_tokens to the model's max completion-token limit."""
    cap = _MODEL_MAX_COMPLETION.get(model)
    if cap and max_tokens > cap:
        return cap
    return max_tokens


def _call(model: str, system: str, user: str, max_tokens: int,
          _truncated: list | None = None) -> str:
    """Call any supported LLM (OpenAI, Gemini, Anthropic, or Claude CLI).

    If *_truncated* is a list, appends True/False to indicate whether the
    response was cut short by the token limit.
    """
    max_tokens = _cap_tokens(model, max_tokens)
    # ── Claude CLI mode: call `claude -p` as subprocess ──────────────────
    if _provider(model) == "claude-cli":
        import subprocess as _sp
        cmd = ["claude", "-p", "--output-format", "text"]
        if system:
            cmd.extend(["--system-prompt", system])
        result = _sp.run(
            cmd, input=user, capture_output=True, text=True,
            timeout=600,
        )
        content = result.stdout.strip()
        if _truncated is not None:
            _truncated.append(False)  # CLI handles its own limits
        if result.returncode != 0 and not content:
            raise RuntimeError(f"claude -p failed (code {result.returncode}): {result.stderr[:500]}")
        return content

    # ── Codex CLI mode: call `codex exec` as subprocess ──────────────────
    # Auth is handled by the `codex` CLI itself (prior `codex login`).
    # We run non-interactively, read-only sandbox, outside a git repo, and
    # capture only the agent's final message via `-o <file>` so we don't
    # have to parse the streaming event log on stdout.
    # The caller's ~/.codex/config.toml default (e.g. gpt-5.2-codex) is
    # often not available on a ChatGPT-plan login, so we override to
    # gpt-5.2 which is broadly available on ChatGPT plans. Set
    # AUTONOTE_CODEX_MODEL to pick a different one — gpt-5.1 requires
    # an API-key codex login, gpt-5.4 / gpt-5.5 / gpt-5.4-mini work on
    # ChatGPT-only accounts.
    if _provider(model) == "codex-cli":
        import subprocess as _sp
        import tempfile as _tf
        import os as _os2
        out_fd, out_file = _tf.mkstemp(prefix="codex_out_", suffix=".txt")
        _os2.close(out_fd)
        try:
            prompt_text = f"{system}\n\n{user}" if system else user
            codex_model = _os2.environ.get("AUTONOTE_CODEX_MODEL", "gpt-5.2")
            cmd = [
                "codex", "exec",
                "-m", codex_model,
                "--skip-git-repo-check",
                "-s", "read-only",
                "-o", out_file,
                "-",  # read prompt from stdin
            ]
            result = _sp.run(
                cmd, input=prompt_text, capture_output=True, text=True,
                timeout=1800,
            )
            try:
                content = Path(out_file).read_text(encoding="utf-8").strip()
            except Exception:
                content = ""
            if _truncated is not None:
                _truncated.append(False)  # CLI handles its own limits
            if result.returncode != 0 and not content:
                err = (result.stderr or result.stdout or "").strip()
                # Surface the real failure at the *tail* of stderr (quota,
                # auth, 400s from the provider). The head is usually just
                # skill-loader warnings and the session preamble.
                tail = err[-600:] if len(err) > 600 else err
                raise RuntimeError(
                    f"codex exec failed (code {result.returncode}): {tail}"
                )
            return content
        finally:
            try:
                Path(out_file).unlink()
            except Exception:
                pass

    client = _get_client_for(model)

    if _provider(model) == "anthropic":
        kwargs: dict = {"model": model, "max_tokens": max_tokens,
                        "messages": [{"role": "user", "content": user}]}
        if system:
            kwargs["system"] = system
        r = client.messages.create(**kwargs)
        if _truncated is not None:
            _truncated.append(r.stop_reason == "max_tokens")
        return r.content[0].text.strip() if r.content else ""

    # OpenAI-compatible (OpenAI + Gemini via OpenAI compat layer)
    msgs = []
    if system:
        msgs.append({"role": "system", "content": system})
    msgs.append({"role": "user", "content": user})

    last_err = None
    for tok in ("max_completion_tokens", "max_tokens"):
        try:
            r = client.chat.completions.create(
                model=model, messages=msgs, **{tok: max_tokens})
            if _truncated is not None:
                reason = getattr(r.choices[0], "finish_reason", None)
                _truncated.append(reason == "length")
            content = r.choices[0].message.content
            return content.strip() if content else ""
        except Exception as e:
            s = str(e)
            if "max_tokens" in s or "max_completion_tokens" in s:
                last_err = e
                continue
            raise
    raise RuntimeError(f"Cannot call {model}: {last_err}")


# ── Slide loading & rendering ─────────────────────────────────────────────────

class SlideInfo:
    __slots__ = ("index", "label", "text", "has_code", "word_count")
    def __init__(self, index: int, label: str, text: str):
        self.index      = index
        self.label      = label
        self.text       = text
        self.has_code   = bool(re.search(
            r"[{};]\s*$|^\s*(int|void|def |class |#include|pthread|malloc)",
            text, re.MULTILINE))
        self.word_count = len(text.split())


def _load_slides(slide_path: Path) -> list[SlideInfo]:
    ext = slide_path.suffix.lower()
    if ext == ".pdf":
        import fitz
        doc = fitz.open(str(slide_path))
        out = []
        for i, page in enumerate(doc):
            text  = page.get_text().strip()
            label = next((ln.strip() for ln in text.splitlines() if ln.strip()), f"Page {i+1}")
            out.append(SlideInfo(i, label[:80], text))
        doc.close()
        return out
    if ext in (".pptx", ".ppt"):
        from pptx import Presentation
        prs = Presentation(str(slide_path))
        out = []
        for i, slide in enumerate(prs.slides):
            parts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        ln = para.text.strip()
                        if ln: parts.append(ln)
            try:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes: parts.append(notes)
            except Exception:
                pass
            text  = "\n".join(parts)
            label = parts[0][:80] if parts else f"Slide {i+1}"
            out.append(SlideInfo(i, label, text))
        return out
    if ext in (".docx", ".doc"):
        from docx import Document
        PAGE_PARA = 15
        doc   = Document(str(slide_path))
        paras = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
        out   = []
        for pi, start in enumerate(range(0, max(len(paras), 1), PAGE_PARA)):
            chunk = paras[start:start + PAGE_PARA]
            text  = "\n".join(chunk)
            label = chunk[0][:80] if chunk else f"Page {pi+1}"
            out.append(SlideInfo(pi, label, text))
        return out
    raise ValueError(f"Unsupported format: {ext}")


def render_slide_images(slide_path: Path, out_dir: Path,
                        indices: list[int] | None = None) -> dict[int, Path]:
    """Render PDF pages to PNG. If indices provided, only render those pages."""
    if slide_path.suffix.lower() != ".pdf":
        return {}
    import fitz
    from PIL import Image as PILImage
    out_dir.mkdir(parents=True, exist_ok=True)
    doc = fitz.open(str(slide_path))
    mat = fitz.Matrix(IMAGE_RENDER_SCALE, IMAGE_RENDER_SCALE)
    mapping: dict[int, Path] = {}
    pages = indices if indices is not None else list(range(len(doc)))
    for i in pages:
        if i >= len(doc):
            continue
        png = out_dir / f"slide_{i+1:03d}.png"
        if not png.exists():
            px  = doc[i].get_pixmap(matrix=mat)
            pil = PILImage.frombytes("RGB", [px.width, px.height], px.samples)
            pil.save(str(png))
            del px, pil
        mapping[i] = png
    doc.close()
    return mapping


# ── Chunk helpers ─────────────────────────────────────────────────────────────

def _clean_artifacts(text: str) -> str:
    """Remove pipeline artifacts that may leak into generated notes."""
    lines = text.splitlines()
    cleaned = []
    for line in lines:
        stripped = line.strip()
        if "terminology or factual errors" in stripped:
            continue
        # Clean section-header artifacts
        line = re.sub(r"##\s*NUS Confidential\s*##", "", line)
        line = re.sub(r"[©(]\s*c?\)?\s*CS\d+", "", line)
        cleaned.append(line)
    return "\n".join(cleaned)


def _ensure_frames_embedded(
    draft: str,
    slides: list,
    img_map: dict,
    img_cache: dict,
    out_dir: Path,
    source: str,
) -> str:
    """Append any frames the LLM forgot to include in its draft.

    The `_build_chunk_prompt` lists every available frame in the prompt
    under "Available images", but DeepSeek V4 (and some other models)
    are inconsistent at actually emitting the `![Frame N](path)` markdown
    even when explicitly asked. This safety net keeps the contract:
    every extracted frame is given a chance to surface in the final
    note. The downstream image-filter pass (`filter_images_pass`) still
    runs and can drop junk frames — we just guarantee they reach that
    pass instead of being silently dropped by the LLM.
    """
    if not img_map:
        return draft
    appended: list[str] = []
    for s in slides:
        if s.index not in img_map:
            continue
        rel = img_map[s.index].relative_to(out_dir)
        rel_str = str(rel).replace("\\", "/")
        # Already cited somewhere in the draft? Skip.
        if rel_str in draft:
            continue
        cache_key = f"page_{s.index}"
        desc = (img_cache.get(cache_key, "") or "").strip()
        # First sentence of the description as the caption (max 140 chars).
        caption = desc[:140]
        for end in ".。！？!?":
            idx = caption.find(end)
            if 25 < idx < len(caption):
                caption = caption[:idx + 1]
                break
        if not caption:
            caption = f"Frame {s.index + 1}" if source == "screenshare" \
                else f"Slide {s.index + 1}"
        prefix = "Frame" if source == "screenshare" else "Slide"
        appended.append(f"![{prefix} {s.index + 1}]({rel_str}) *({caption})*")
    if not appended:
        return draft
    return draft.rstrip() + "\n\n" + "\n\n".join(appended) + "\n"


_BAD_LABEL = re.compile(
    r"^\s*(\d+|[A-Z]{2,4}\d{4}[\s\-].*|CS\d+.*|AY\d+.*|\[.*\]|"
    r".*NUS Confidential.*|.*©\s*CS\d+.*|\(c\)\s*CS\d+.*|Page\s+\d+)\s*$",
    re.IGNORECASE,
)

def _dedup_slides(slides: list[SlideInfo], threshold: float = 0.85) -> list[SlideInfo]:
    """Remove near-duplicate slides by comparing text content.

    Groups consecutive slides whose text Jaccard similarity exceeds
    *threshold* and keeps only the one with the most text (typically
    the "fully revealed" version of an incremental slide).  Also deduplicates
    non-consecutive slides that are very similar (threshold 0.9).
    """
    if len(slides) <= 1:
        return slides

    def _words(s: SlideInfo) -> set[str]:
        return set(s.text.lower().split())

    # Pass 1: merge consecutive near-duplicates
    kept: list[SlideInfo] = [slides[0]]
    for s in slides[1:]:
        wa, wb = _words(kept[-1]), _words(s)
        union = wa | wb
        if union and len(wa & wb) / len(union) >= threshold:
            # Keep the one with more content
            if s.word_count > kept[-1].word_count:
                kept[-1] = s
        else:
            kept.append(s)

    # Pass 2: remove non-consecutive near-duplicates (very high threshold)
    final: list[SlideInfo] = []
    seen_texts: list[set[str]] = []
    for s in kept:
        ws = _words(s)
        is_dup = False
        for prev_ws in seen_texts:
            union = ws | prev_ws
            if union and len(ws & prev_ws) / len(union) >= 0.9:
                is_dup = True
                break
        if not is_dup:
            final.append(s)
            seen_texts.append(ws)

    return final


def _chunk_title(slides_in_chunk: list[SlideInfo]) -> str:
    """Pick a representative title for a chunk of slides.

    Prefer short, meaningful slide labels (section headers).
    Skip labels that are pure numbers, course codes, or bracket tags.
    """
    def _is_good(label: str) -> bool:
        if not label or len(label) < 4:
            return False
        if _BAD_LABEL.match(label):
            return False
        # skip labels that are just digits or single tokens that look like slide numbers
        if re.match(r"^\d+$", label.strip()):
            return False
        return True

    # Prefer short title-like labels (section headers are typically ≤8 words)
    for s in slides_in_chunk:
        words = len(s.label.split())
        if 1 <= words <= 8 and _is_good(s.label):
            return s.label
    # Fall back to first slide with a good label
    for s in slides_in_chunk:
        if _is_good(s.label):
            return s.label
    # Last resort: use slide range as title instead of a bad label
    first = slides_in_chunk[0].index + 1
    last  = slides_in_chunk[-1].index + 1
    return f"Slides {first}–{last}" if first != last else f"Slide {first}"


def _build_chunk_prompt(
    slides: list[SlideInfo],
    compact_by_idx: dict[int, dict],
    img_cache: dict,
    img_map: dict[int, Path],
    out_dir: Path,
    course_name: str,
    lec_num: int,
    lec_title: str,
    chunk_idx: int,
    chunk_title: str,
    detail: int,
    has_transcript: bool,
    source: str = "slides",
) -> str:
    # Slide outline: number, title, first line of text (skip page numbers/headers)
    outline_lines = []
    for s in slides:
        if source == "screenshare":
            # For screen share, the "slide" is a video frame — no text content,
            # but the transcript carries all the information
            outline_lines.append(f"  Frame {s.index+1}: (screen capture at this point)")
        else:
            # Extract meaningful text (skip short lines like "10", "[ CS3210 ]")
            meaningful = [ln for ln in s.text.splitlines()
                          if len(ln.strip()) > 8 and not re.match(r"^\d+$", ln.strip())]
            snippet = " · ".join(meaningful[:3])[:120]
            outline_lines.append(f"  Slide {s.index+1}: 「{s.label}」  {snippet}")
    slide_outline = "\n".join(outline_lines)

    # Image hints — only include images with visual content (diagrams, code,
    # formulas, charts).  Pure-text slides/frames add no value as images since
    # the text is already in the transcript or slide outline.
    img_hints_lines = []
    for s in slides:
        if s.index not in img_map:
            continue
        rel = img_map[s.index].relative_to(out_dir)

        # Build context: what the lecturer was saying when this slide was shown
        cs = compact_by_idx.get(s.index)
        transcript_ctx = ""
        if cs and cs.get("transcript", "").strip():
            transcript_ctx = cs["transcript"][:150].replace("\n", " ")

        if source == "screenshare":
            cache_key = f"page_{s.index}"
            desc = img_cache.get(cache_key, "").strip()
            brief = desc[:220] if desc else (s.text[:120].replace("\n", " ") if s.text.strip() else "")
            ctx = f" [context: {transcript_ctx}]" if transcript_ctx else ""
            img_hints_lines.append(f"  Frame {s.index+1}: `{rel}` — {brief}{ctx}")
        else:
            cache_key = f"page_{s.index}"
            desc = img_cache.get(cache_key, "")
            if desc or s.word_count < 80 or s.has_code:
                note = desc if desc else ("has code" if s.has_code else s.label)
                ctx = f" [context: {transcript_ctx}]" if transcript_ctx else ""
                img_hints_lines.append(f"  Slide {s.index+1}: `{rel}` — {note}{ctx}")
    image_hints = "\n".join(img_hints_lines) or "  (no images for this segment)"

    if has_transcript:
        # Transcript block: [MM:SS Slide N「Title」] transcript (full text).
        # For slides without transcript, fall back to the cached image
        # description so the LLM still has content to write notes about them.
        transcript_lines = []
        for s in slides:
            cs = compact_by_idx.get(s.index)
            tx = cs.get("transcript", "").strip() if cs else ""
            if tx:
                mm = int(cs["start"] // 60)
                ss = int(cs["start"] % 60)
                transcript_lines.append(
                    f"[{mm:02d}:{ss:02d} Slide {s.index+1}「{s.label}」]\n{tx}")
            else:
                # No transcript — use image description if available
                desc = img_cache.get(f"page_{s.index}", "").strip() if source == "screenshare" else ""
                if desc and len(desc) > 30:
                    transcript_lines.append(
                        f"[Slide {s.index+1}「{s.label}」 — no audio, visual only]\n"
                        f"(Visual content: {desc[:300]})")
        transcript_block = "\n\n".join(transcript_lines) or _P("no_transcript")

        prompt = _P("chunk").format(
            course_name=course_name, lec_num=lec_num, lec_title=lec_title,
            slide_outline=slide_outline, transcript_block=transcript_block,
            image_hints=image_hints, chunk_idx=chunk_idx, chunk_title=chunk_title,
            detail=detail, detail_instruction=_detail_instr(detail),
        )
        # Truncate if total prompt exceeds limit
        if len(prompt) > MAX_NOTE_CHARS:
            prompt = prompt[:MAX_NOTE_CHARS] + "\n\n[...transcript truncated due to length...]"
        return prompt
    else:
        return _P("slide_only").format(
            course_name=course_name, lec_num=lec_num, lec_title=lec_title,
            slide_outline=slide_outline, image_hints=image_hints,
            chunk_idx=chunk_idx, chunk_title=chunk_title,
            detail=detail, detail_instruction=_detail_instr(detail),
        )


# ── Section-by-section generation ────────────────────────────────────────────

def _section_path(sections_dir: Path, ld: "LectureData", ci: int) -> Path:
    # Keyed by `dir_key` (L{num}_{slug}[_F{idx}]) so the cache survives
    # caption-list reshuffles — a section generated for one lecture can
    # never be served to a different lecture that happens to land on the
    # same sequential `num` on a later run.
    return sections_dir / f"{ld.dir_key}_S{ci:02d}.md"


def generate_section(
    lec_num: int,
    lec_title: str,
    course_name: str,
    chunk: list[SlideInfo],
    ci: int,
    ld: "LectureData",
    out_dir: Path,
    sections_dir: Path,
    detail: int,
    has_transcript: bool,
    bar: tqdm | None = None,
    force: bool = False,
) -> str:
    """Generate (or load from cache) one section and save it to sections_dir."""
    sec_file = _section_path(sections_dir, ld, ci)

    if not force and sec_file.exists() and sec_file.stat().st_size > 500:
        if bar:
            fi = f"F{ld.file_idx} " if ld.file_idx > 1 else ""
            bar.set_postfix_str(f"L{lec_num}{fi}§{ci} cached")
        return sec_file.read_text(encoding="utf-8"), False

    chunk_title = _chunk_title(chunk)
    fi_tag = f"F{ld.file_idx} " if ld.file_idx > 1 else ""
    if bar:
        bar.set_postfix_str(f"L{lec_num}{fi_tag}§{ci}/{chunk_title[:22]} generating")
    tqdm.write(f"  → L{lec_num}{fi_tag}§{ci} [{chunk_title[:40]}]  slides {chunk[0].index+1}–{chunk[-1].index+1}  calling {NOTE_MODEL}…")

    # Render only this chunk's slide images (lazy, avoids OOM)
    img_map = ld.render_chunk_images([s.index for s in chunk])

    user = _build_chunk_prompt(
        slides=chunk,
        compact_by_idx=ld.compact_by_idx,
        img_cache=ld.img_cache,
        img_map=img_map,
        out_dir=out_dir,
        course_name=course_name,
        lec_num=lec_num,
        lec_title=lec_title,
        chunk_idx=ci,
        chunk_title=chunk_title,
        detail=detail,
        has_transcript=has_transcript,
        source=ld.source,
    )

    import time as _time
    _t0 = _time.monotonic()
    _trunc_flag: list[bool] = []
    draft = _call(NOTE_MODEL, _P("system"), user, _max_tokens(detail),
                  _truncated=_trunc_flag)
    was_truncated = bool(_trunc_flag and _trunc_flag[0])
    # Belt-and-suspenders: when the provider returns finish_reason='stop'
    # but the content was actually cut (DeepSeek does this), the shape
    # heuristic catches it.
    if not was_truncated and _looks_truncated(draft):
        was_truncated = True
        tqdm.write(f"     [warn] Draft looks truncated despite stop reason — "
                   f"flagging for retry")
    tqdm.write(f"     ✓ {len(draft):,} chars  ({_time.monotonic()-_t0:.0f}s)"
               + ("  [TRUNCATED]" if was_truncated else ""))

    if not draft or len(draft.strip()) < 100:
        tqdm.write(f"  [warn] Empty or too-short draft for L{lec_num} §{ci} — not caching")
        heading = f"### {lec_num}.{ci} {_chunk_title(chunk)}"
        return f"{heading}\n\n*(Section could not be generated — re-run with force to retry.)*", True

    # Strip pipeline artifacts that may have leaked into the draft
    draft = _clean_artifacts(draft)

    # Frame-completeness safety net: DeepSeek V4 (and other models under
    # certain phrasings) is inconsistent at embedding `![Frame N](...)`
    # markdown even when the prompt asks for it — empirically the LLM
    # keeps anywhere from 0–100% of the available frames per chunk.
    # Auto-append any frame in img_map that didn't make it into the draft
    # so we never silently lose the visual content the user asked us to
    # extract. The downstream image-filter pass still drops junk frames.
    draft = _ensure_frames_embedded(
        draft, chunk, img_map, ld.img_cache, out_dir, ld.source,
    )

    # Translate to target language if not English
    if NOTE_LANGUAGE != "en" and draft:
        lang = _LANG_NAMES.get(NOTE_LANGUAGE, NOTE_LANGUAGE)
        tqdm.write(f"     translating to {lang}…")
        _tt = _time.monotonic()
        try:
            draft = _translate(draft, lang)
            tqdm.write(f"     ✓ translated ({_time.monotonic()-_tt:.0f}s)")
        except Exception as _te:
            tqdm.write(f"     [warn] Translation failed "
                       f"({type(_te).__name__}: {str(_te)[:120]}) — "
                       f"keeping draft in source language")

    heading = f"### {lec_num}.{ci} {chunk_title}"
    content = f"{heading}\n\n{draft}"
    if was_truncated:
        tqdm.write(f"  [warn] Section L{lec_num} §{ci} was truncated by token limit — not caching")
    else:
        sec_file.write_text(content, encoding="utf-8")
    return content, True


def generate_lecture(
    lec_num: int,
    lec_title: str,
    course_name: str,
    ld: "LectureData",
    out_dir: Path,
    sections_dir: Path,
    detail: int,
    fmt: str,
    bar: tqdm | None = None,
    force: bool = False,
) -> str:
    """Generate all sections for one lecture, saving each to sections_dir.
    Returns (merged_text, any_fresh) where any_fresh is True if at least
    one section was freshly generated (not loaded from cache)."""
    has_transcript = bool(ld.compact_by_idx)
    any_fresh = False

    # Deduplicate near-identical slides (e.g. incremental reveals, repeated
    # content across pre-lecture/main slides).  Keep the slide with the most
    # text content from each group of similar slides.
    slides = _dedup_slides(ld.slides)
    if len(slides) < len(ld.slides):
        tqdm.write(f"  Deduplicated: {len(ld.slides)} → {len(slides)} slides "
                   f"(removed near-duplicates)")

    chunks = [slides[i:i+CHAPTER_SIZE]
              for i in range(0, len(slides), CHAPTER_SIZE)]

    # Pre-render all slide images before parallel generation so threads
    # don't race on the shared ld.img_map dict and filesystem writes.
    all_indices = [s.index for s in slides]
    ld.render_chunk_images(all_indices)

    # Parallel section generation.  Each section's generate + translate is
    # independent.  claude-cli spawns subprocesses (I/O bound, not CPU bound),
    # so we can run many concurrently just like API calls.
    PARALLEL_SECTIONS = 6
    tqdm.write(f"  [{len(chunks)} chunks, {PARALLEL_SECTIONS} workers → "
               f"{'parallel' if len(chunks) > 1 else 'sequential'}]")
    if len(chunks) <= 1 or PARALLEL_SECTIONS <= 1:
        # Sequential fallback for single-section lectures
        parts: list[str] = []
        for ci, chunk in enumerate(chunks, start=1):
            content, fresh = generate_section(
                lec_num=lec_num, lec_title=lec_title, course_name=course_name,
                chunk=chunk, ci=ci, ld=ld, out_dir=out_dir,
                sections_dir=sections_dir, detail=detail,
                has_transcript=has_transcript, bar=bar, force=force,
            )
            parts.append(content)
            any_fresh = any_fresh or fresh
            if bar:
                bar.update(1)
    else:
        from concurrent.futures import ThreadPoolExecutor, as_completed
        import threading
        _bar_lock = threading.Lock()

        def _gen(ci_chunk):
            ci, chunk = ci_chunk
            return ci, generate_section(
                lec_num=lec_num, lec_title=lec_title, course_name=course_name,
                chunk=chunk, ci=ci, ld=ld, out_dir=out_dir,
                sections_dir=sections_dir, detail=detail,
                has_transcript=has_transcript, bar=None, force=force,
            )

        results: dict[int, tuple[str, bool]] = {}
        with ThreadPoolExecutor(max_workers=PARALLEL_SECTIONS) as pool:
            futures = {pool.submit(_gen, (ci, chunk)): ci
                       for ci, chunk in enumerate(chunks, start=1)}
            for fut in as_completed(futures):
                ci_ret, (content, fresh) = fut.result()
                results[ci_ret] = (content, fresh)
                if bar:
                    with _bar_lock:
                        bar.update(1)

        # Reassemble in order
        parts = []
        for ci in sorted(results.keys()):
            content, fresh = results[ci]
            parts.append(content)
            any_fresh = any_fresh or fresh

    return "\n\n".join(parts), any_fresh


# ── Self-scoring ──────────────────────────────────────────────────────────────

def _key_terms(text: str, n: int = 12) -> list[str]:
    tokens = re.findall(r"\b[A-Z][a-zA-Z_]{2,}\b|\b[A-Z_]{3,}\b", text)
    cs = re.findall(
        r"\b(mutex|semaphore|thread|process|deadlock|fork|kernel|scheduler|"
        r"critical.section|race.condition|barrier|starvation|spinlock|heap|"
        r"socket|packet|routing|subnet|protocol|checksum|TCP|UDP|ACK|NAK|"
        r"Viterbi|pthread|register|pipeline|cache|interrupt|syscall)\b",
        text, re.IGNORECASE,
    )
    combined = [t.lower() for t in tokens + cs]
    seen: set[str] = set()
    result: list[str] = []
    for t in sorted(set(combined), key=len, reverse=True):
        if t not in seen:
            seen.add(t); result.append(t)
        if len(result) >= n: break
    return result


def self_score(all_slides: list[SlideInfo], full_notes: str,
               all_compact: list[dict]) -> dict:
    notes_lower = full_notes.lower()

    # Coverage based on word count vs expected
    word_count      = len(full_notes.split())
    expected_words  = len(all_slides) * MIN_NOTE_WORDS_PER_SLIDE
    coverage_score  = min(word_count / expected_words, 1.0) * 10 if expected_words else 10

    # Terminology
    term_hits, term_total = 0, 0
    for s in all_slides:
        terms = _key_terms(s.text)
        if terms:
            term_total += len(terms)
            term_hits  += sum(1 for t in terms if t in notes_lower)
    terminology_score = (term_hits / term_total * 10) if term_total else 10

    # Callouts
    compact_map = {(c["slide"] - 1): c for c in all_compact}
    callouts_needed  = 0
    callouts_written = len(re.findall(r"\[!IMPORTANT\]", full_notes))
    for s in all_slides:
        cs = compact_map.get(s.index)
        if cs:
            tx = cs.get("transcript", "").lower()
            if any(w in tx for w in ("important", "remember", "exam", "key point",
                                      "critical", "note that", "this will be")):
                callouts_needed += 1
    callout_score = (min(callouts_written, callouts_needed) / callouts_needed * 10) \
                    if callouts_needed else 10

    # Code blocks
    code_needed  = sum(1 for s in all_slides if s.has_code)
    code_written = len(re.findall(r"```\w+", full_notes))
    code_score   = (min(code_written, code_needed) / code_needed * 10) if code_needed else 10

    overall = (SCORE_WEIGHTS["coverage"]    * coverage_score
             + SCORE_WEIGHTS["terminology"] * terminology_score
             + SCORE_WEIGHTS["callouts"]    * callout_score
             + SCORE_WEIGHTS["code_blocks"] * code_score)

    return {
        "coverage":    round(coverage_score,    1),
        "terminology": round(terminology_score, 1),
        "callouts":    round(callout_score,      1),
        "code_blocks": round(code_score,         1),
        "overall":     round(overall,            2),
        "stats": {
            "note_words":       word_count,
            "expected_words":   expected_words,
            "term_hits":        term_hits,
            "term_total":       term_total,
            "callouts_written": callouts_written,
            "callouts_needed":  callouts_needed,
            "code_blocks":      code_written,
            "code_slides":      code_needed,
        },
    }


def _print_score(scores: dict, label: str) -> None:
    st = scores.get("stats", {})
    tqdm.write(f"\n  ┌──────────────────────────────────────────────┐")
    tqdm.write(f"  │ Score: {label[:38]:38s}│")
    tqdm.write(f"  ├──────────────────────────────────────────────┤")
    tqdm.write(f"  │ Coverage   {scores['coverage']:4.1f}/10  "
               f"({st.get('note_words','?')} / ~{st.get('expected_words','?')} words) │")
    tqdm.write(f"  │ Terminology{scores['terminology']:4.1f}/10  "
               f"({st.get('term_hits','?')}/{st.get('term_total','?')} terms)       │")
    tqdm.write(f"  │ Callouts   {scores['callouts']:4.1f}/10  "
               f"({st.get('callouts_written','?')}/{st.get('callouts_needed','?')} callouts)     │")
    tqdm.write(f"  │ Code blocks{scores['code_blocks']:4.1f}/10  "
               f"({st.get('code_blocks','?')}/{st.get('code_slides','?')} code slides)   │")
    tqdm.write(f"  │ Overall    {scores['overall']:4.2f}/10                              │")
    tqdm.write(f"  └──────────────────────────────────────────────┘")


# ── LectureData ───────────────────────────────────────────────────────────────

class LectureData:
    def __init__(self, num: int, slide_path: Path, alignment_path: Path | None,
                 file_idx: int = 1, source: str = "slides",
                 frame_dir: Path | None = None):
        self.num            = num
        self.file_idx       = file_idx
        self.slide_path     = slide_path
        self.alignment_path = alignment_path
        self.source         = source       # "slides" or "screenshare"
        self.frame_dir      = frame_dir    # directory of extracted frame PNGs
        self.slides:         list[SlideInfo] = []
        self.compact:        dict            = {}
        self.compact_slides: list[dict]      = []
        self.img_cache:      dict            = {}
        self.img_map:        dict[int, Path] = {}

    @property
    def slug(self) -> str:
        # Stable, filesystem-safe identifier derived from the lecture's
        # source. Used to disambiguate per-lecture image directories and
        # section cache files so different lectures never share paths —
        # previously `num` (a sequential index over the captions list)
        # shifted when videos were added, causing e.g. 4_17's frames to
        # overwrite 4_10's inside images/L11/ and leaving stale captions
        # pointing at the wrong images.
        if self.source == "screenshare" and self.frame_dir is not None:
            base = self.frame_dir.name
        else:
            base = self.slide_path.stem
        safe = re.sub(r"[^A-Za-z0-9_-]+", "_", base).strip("_")[:48]
        return safe or "lec"

    @property
    def dir_key(self) -> str:
        # The per-lecture path component for both images/ and the section
        # cache: `L{num:02d}_{slug}` (plus `_F{file_idx:02d}` when a single
        # lecture has more than one slide file).
        base = f"L{self.num:02d}_{self.slug}"
        if self.file_idx > 1:
            base = f"{base}_F{self.file_idx:02d}"
        return base

    def load(self, out_dir: Path) -> None:
        if self.source == "screenshare" and self.frame_dir:
            self._load_from_frames(out_dir)
        elif self.source == "transcript_only":
            self._load_from_transcript(out_dir)
        else:
            self.slides = _load_slides(self.slide_path)
        if self.alignment_path and self.alignment_path.exists() \
                and self.source != "transcript_only":
            self.compact = alignment_parser.parse(self.alignment_path)
            self.compact_slides = self.compact.get("slides", [])
        if self.source == "slides":
            cache_f = self.slide_path.parent / f"{self.slide_path.name}.image_cache.json"
            if cache_f.exists():
                with open(cache_f) as f:
                    self.img_cache = json.load(f)
        elif self.source == "screenshare":
            # For screenshare: load the frame description cache created by
            # frame_extractor._describe_frames()
            cache_f = self.frame_dir / "image_cache.json"
            if cache_f.exists():
                with open(cache_f) as f:
                    self.img_cache = json.load(f)
        # img_map is populated lazily per chunk to avoid rendering all slides at once
        self._out_dir = out_dir

    def _load_from_frames(self, out_dir: Path) -> None:
        """Load 'slides' from extracted video frames instead of PDF/PPTX."""
        frame_pngs = sorted(self.frame_dir.glob("frame_*.png"))
        self.slides = []
        for i, png in enumerate(frame_pngs):
            # For screen share frames, the text content is minimal (just the label).
            # The actual content comes from the transcript in the alignment.
            label = f"Frame {i + 1}"
            text = label  # placeholder; transcript provides the real content
            self.slides.append(SlideInfo(i, label, text))
            # Pre-populate img_map since frames are already extracted
            self.img_map[i] = png

    def _load_from_transcript(self, out_dir: Path) -> None:
        """Build pseudo-slides from a caption when no slide file or screenshare
        frames are available. Each pseudo-slide covers ~CHUNK_SEC of audio, so
        the per-chunk LLM call receives a coherent transcript segment."""
        CHUNK_SEC = 180.0   # ~3 minutes per pseudo-slide
        try:
            with open(self.slide_path, encoding="utf-8") as f:
                cap = json.load(f)
        except Exception:
            self.slides = [SlideInfo(0, "Lecture", "")]
            return
        segs = cap.get("segments", [])
        if not segs:
            self.slides = [SlideInfo(0, "Lecture", "")]
            return

        # Chunk segments into pseudo-slides by time.
        chunks: list[list[dict]] = []
        cur: list[dict] = []
        cur_start = segs[0].get("start", 0.0)
        for s in segs:
            if s.get("start", 0.0) - cur_start >= CHUNK_SEC and cur:
                chunks.append(cur)
                cur = []
                cur_start = s.get("start", 0.0)
            cur.append(s)
        if cur:
            chunks.append(cur)

        self.slides = []
        compact_slides: list[dict] = []
        for i, chunk in enumerate(chunks):
            text = " ".join(s.get("text", "").strip() for s in chunk).strip()
            label = f"Segment {i + 1}"
            self.slides.append(SlideInfo(i, label, text))
            compact_slides.append({
                "slide":      i + 1,
                "start":      chunk[0].get("start", 0.0),
                "end":        chunk[-1].get("end", 0.0),
                "transcript": text,
            })
        self.compact_slides = compact_slides
        self.compact = {
            "lecture":      self.slide_path.stem,
            "slides":       compact_slides,
            "duration":     cap.get("duration", 0.0),
            "source":       "transcript_only",
        }

    def render_chunk_images(self, slide_indices: list[int]) -> dict[int, Path]:
        """Render only the slides in this chunk into images/{dir_key}/, cache results."""
        if self.source == "transcript_only":
            # No slide file and no frames to render.
            return {}
        if self.source == "screenshare":
            # For screen share, frames are already extracted — copy them
            # into the notes images/ directory using frame_NNN naming.
            # `dir_key` embeds the lecture slug so unrelated lectures
            # that happen to share the same `num` don't clobber each
            # other's frames under a common images/L{num}/ path.
            img_dir = self._out_dir / "images" / self.dir_key
            img_dir.mkdir(parents=True, exist_ok=True)
            import shutil
            for i in slide_indices:
                if i in self.img_map:
                    src = self.img_map[i]
                    dst = img_dir / f"frame_{i + 1:03d}.png"
                    # Always refresh: frames may have been re-extracted with
                    # different content under the same number after a dedup
                    # or junk-filter pass.
                    if (not dst.exists() or
                            src.stat().st_mtime > dst.stat().st_mtime or
                            src.stat().st_size != dst.stat().st_size):
                        shutil.copy2(str(src), str(dst))
                    self.img_map[i] = dst
            return {i: self.img_map[i] for i in slide_indices if i in self.img_map}

        img_dir = self._out_dir / "images" / self.dir_key
        needed = [i for i in slide_indices if i not in self.img_map]
        if needed:
            new = render_slide_images(self.slide_path, img_dir, needed)
            self.img_map.update(new)
        return {i: self.img_map[i] for i in slide_indices if i in self.img_map}

    @property
    def title(self) -> str:
        # Prefer alignment lecture title (e.g. "CS3210 e-Lecture on Processes and Threads")
        if self.compact.get("lecture"):
            t = self.compact["lecture"]
            # Strip course prefix like "CS3210 e-Lecture on "
            t = re.sub(r"^CS\d+\s+e-Lecture\s+on\s+", "", t)
            # Strip trailing speaker credit "(by ...)"
            t = re.sub(r"\s*\(by .+\)$", "", t)
            return t.strip()
        if self.source == "screenshare" and self.frame_dir:
            stem = self.frame_dir.name
            return stem.replace("-", " ").replace("_", " ")
        if self.source == "transcript_only":
            stem = self.slide_path.stem
            return stem.replace("-", " ").replace("_", " ")
        # Fall back to filename stem, removing "L02-" prefix
        stem = re.sub(r"^[Ll]\d+[-_\s]+", "", self.slide_path.stem)
        return stem.replace("-", " ").replace("_", " ") or self.slide_path.stem

    @property
    def compact_by_idx(self) -> dict[int, dict]:
        return {(c["slide"] - 1): c for c in self.compact_slides}


# ── Merge sections into final note ────────────────────────────────────────────

def merge_sections(
    course_name: str,
    lectures: list["LectureData"],
    sections_dir: Path,
    out_path: Path,
    all_slides: list[SlideInfo],
    all_compact: list[dict],
    run_image_filter: bool = True,
) -> tuple[Path, dict]:
    """Read all saved section files and merge into one final Markdown note."""
    from itertools import groupby

    note_sections: list[str] = []

    for lec_num, ld_iter in groupby(lectures, key=lambda x: x.num):
        ld_group   = list(ld_iter)
        multi_file = len(ld_group) > 1
        ld0 = ld_group[0]
        video_name = ld0.compact.get("lecture", ld0.slide_path.stem)
        slide_name = ld0.slide_path.name if ld0.source != "screenshare" else "(screen recording)"
        lec_heading = (f"## Lecture {lec_num} — {ld0.title}\n\n"
                       f"*Video: {video_name} | Slides: {slide_name}*")
        lec_parts: list[str] = []

        for ld in ld_group:
            n_chunks   = max(1, (len(ld.slides) + CHAPTER_SIZE - 1) // CHAPTER_SIZE)
            file_parts: list[str] = []
            for ci in range(1, n_chunks + 1):
                sec_file = _section_path(sections_dir, ld, ci)
                if sec_file.exists() and sec_file.stat().st_size > 500:
                    file_parts.append(sec_file.read_text(encoding="utf-8"))
                else:
                    fi = f" F{ld.file_idx}" if multi_file else ""
                    tqdm.write(f"  [warn] Missing section L{ld.num}{fi} §{ci} — skipping")

            if file_parts:
                if multi_file:
                    part_heading = f"### Part {ld.file_idx}: {ld.slide_path.stem}"
                    lec_parts.append(f"{part_heading}\n\n" + "\n\n".join(file_parts))
                else:
                    lec_parts.extend(file_parts)

        if lec_parts:
            note_sections.append(f"{lec_heading}\n\n" + "\n\n".join(lec_parts))

    # Exam notes — generated from the merged content
    tqdm.write("  Generating exam notes…")
    summary = "\n\n---\n\n".join(note_sections)[:8000]
    exam_section_file = sections_dir / "exam_notes.md"
    if exam_section_file.exists() and exam_section_file.stat().st_size > 50:
        exam_md = exam_section_file.read_text(encoding="utf-8")
    else:
        exam_md = _call(NOTE_MODEL, _P("system"),
                        _P("exam").format(course_name=course_name, summary=summary),
                        4000)
        exam_section_file.write_text(exam_md, encoding="utf-8")
    note_sections.append(exam_md)

    now = datetime.now().strftime("%Y-%m-%dT%H:%M:%S+08:00")
    front = (f"---\ntitle: {course_name}\ndate: {now}\n"
             f"description: Lecture notes generated by auto_note\n"
             f"categories:\n    - tech\n---\n\n")
    full_notes = front + f"# {course_name} Notes\n\n" + "\n\n------\n\n".join(note_sections)

    # ── Image filter agent pass (only for newly generated content) ─────────────
    if run_image_filter:
        tqdm.write("  Running image filter pass…")
        full_notes, _, _ = filter_images_pass(full_notes, out_path.parent, lectures)
    else:
        tqdm.write("  Skipping image filter (all sections from cache).")

    out_path.write_text(full_notes, encoding="utf-8")
    tqdm.write(f"\n  Merged → {out_path}  ({len(full_notes):,} chars)")

    scores = {}
    if SHOW_SCORE:
        scores = self_score(all_slides, full_notes, all_compact)
        _print_score(scores, out_path.name)
        score_path = out_path.with_suffix(".score.json")
        with open(score_path, "w") as f:
            json.dump(scores, f, indent=2)

    return out_path, scores


# ── Generate full course notes ────────────────────────────────────────────────

def generate_course_notes(
    course_name: str,
    lectures: list[LectureData],
    out_path: Path,
    detail: int = DETAIL_LEVEL,
    fmt: str = OUTPUT_FORMAT,
    force: bool = False,
) -> tuple[Path, dict]:

    out_dir = out_path.parent
    out_dir.mkdir(parents=True, exist_ok=True)

    # sections/ sub-directory: each chunk saved as L{N}_S{ci}.md
    sections_dir = out_dir / "sections"
    sections_dir.mkdir(exist_ok=True)
    # Write language marker so future runs detect language changes
    (sections_dir / ".language").write_text(NOTE_LANGUAGE, encoding="utf-8")

    # Load all lectures first to get accurate chunk counts
    tqdm.write("  Loading lecture data…")
    for ld in lectures:
        ld.load(out_dir)

    total_chunks = sum(
        max(1, (len(ld.slides) + CHAPTER_SIZE - 1) // CHAPTER_SIZE)
        for ld in lectures
    )

    all_slides:  list[SlideInfo] = []
    all_compact: list[dict]      = []

    bar = tqdm(total=total_chunks, desc=f"{course_name}", unit="section",
               bar_format="{l_bar}{bar}| {n_fmt}/{total_fmt} sections [{elapsed}<{remaining}]")

    # ── Phase 1: generate each section independently ──────────────────────────
    any_new = False
    for ld in lectures:
        all_slides.extend(ld.slides)
        all_compact.extend(ld.compact_slides)
        _, fresh = generate_lecture(
            lec_num=ld.num,
            lec_title=ld.title,
            course_name=course_name,
            ld=ld,
            out_dir=out_dir,
            sections_dir=sections_dir,
            detail=detail,
            fmt=fmt,
            bar=bar,
            force=force,
        )
        any_new = any_new or fresh

    bar.close()

    # ── Phase 2: merge all sections into one file ─────────────────────────────
    tqdm.write("\n  Merging sections…")
    return merge_sections(
        course_name=course_name,
        lectures=lectures,
        sections_dir=sections_dir,
        out_path=out_path,
        all_slides=all_slides,
        all_compact=all_compact,
        run_image_filter=any_new,
    )


# ── Per-video note generation ─────────────────────────────────────────────────

def generate_per_video_notes(
    course_name: str,
    lectures: list[LectureData],
    out_dir: Path,
    detail: int = DETAIL_LEVEL,
    fmt: str = OUTPUT_FORMAT,
    force: bool = False,
) -> list[tuple[Path, dict]]:
    """Generate one note file per lecture number instead of a single merged file.

    Lectures sharing the same lecture number (multi-part) are merged into
    one file: <course>_L<N>_notes.md.
    Returns a list of (path, scores) tuples for each generated note.
    """
    from itertools import groupby

    out_dir.mkdir(parents=True, exist_ok=True)
    sections_dir = out_dir / "sections"
    sections_dir.mkdir(exist_ok=True)
    (sections_dir / ".language").write_text(NOTE_LANGUAGE, encoding="utf-8")

    tqdm.write("  Loading lecture data…")
    for ld in lectures:
        ld.load(out_dir)

    results: list[tuple[Path, dict]] = []

    for ld in lectures:
        lec_num = ld.num
        lec_title = ld.title
        # Name the note file after the video (caption stem), not the lecture number
        video_name = ld.compact.get("lecture", ld.slide_path.stem)
        # Sanitize for filename
        safe_name = re.sub(r'[\\/*?:"<>|]', '_', video_name).strip()
        ext = ".mdx" if fmt == "mdx" else ".md"
        note_path = out_dir / f"{safe_name}_notes{ext}"

        tqdm.write(f"\n  ═══ {video_name} ═══")

        n_chunks = max(1, (len(ld.slides) + CHAPTER_SIZE - 1) // CHAPTER_SIZE)
        bar = tqdm(total=n_chunks, desc=safe_name[:30],
                   unit="section",
                   bar_format="{l_bar}{bar}| {n_fmt}/{total_fmt} [{elapsed}<{remaining}]")

        text, fresh = generate_lecture(
            lec_num=lec_num, lec_title=lec_title,
            course_name=course_name, ld=ld, out_dir=out_dir,
            sections_dir=sections_dir, detail=detail, fmt=fmt,
            bar=bar, force=force,
        )
        bar.close()

        # Source metadata
        if ld.source == "screenshare":
            slide_file = "(screen recording frames)"
        elif ld.source == "transcript_only":
            slide_file = "(transcript only — no slides matched)"
        else:
            slide_file = ld.slide_path.name
        source_info = f"- **Video**: {video_name}\n- **Slides**: {slide_file}\n\n---\n\n"

        from datetime import datetime
        now = datetime.now().strftime("%Y-%m-%dT%H:%M:%S+08:00")
        front = (f"---\ntitle: {course_name} — {video_name}\ndate: {now}\n"
                 f"description: Lecture notes generated by auto_note\n"
                 f"categories:\n    - tech\n---\n\n")
        heading = f"# {course_name} — {video_name}\n\n"
        full_notes = front + heading + source_info + text

        if fresh:
            tqdm.write("  Running image filter pass…")
            full_notes, _, _ = filter_images_pass(full_notes, out_dir, [ld])
        else:
            tqdm.write("  Skipping image filter (all sections from cache).")

        note_path.write_text(full_notes, encoding="utf-8")
        tqdm.write(f"  Saved → {note_path}  ({len(full_notes):,} chars)")

        results.append((note_path, {}))

    return results


# ── Iteration ─────────────────────────────────────────────────────────────────

def generate_with_iteration(
    course_name: str,
    lectures: list[LectureData],
    out_path: Path,
    fmt: str = OUTPUT_FORMAT,
    max_rounds: int = 3,
) -> Path:
    detail     = DETAIL_LEVEL
    best_path  = out_path
    best_score = 0.0

    for rnd in range(1, max_rounds + 1):
        print(f"\n{'#'*60}")
        print(f"# ROUND {rnd}/{max_rounds}   detail={detail}   target={QUALITY_TARGET}")
        print(f"{'#'*60}")

        path, scores = generate_course_notes(
            course_name, lectures, out_path,
            detail=detail, fmt=fmt,
            force=(rnd > 1),   # re-generate sections on subsequent rounds
        )
        overall = scores["overall"]
        if overall > best_score:
            best_score = overall
            best_path  = path

        if overall >= QUALITY_TARGET:
            tqdm.write(f"\n  ✓ Target {QUALITY_TARGET} reached (score={overall:.2f})")
            break
        if rnd < max_rounds:
            tqdm.write(f"\n  Score {overall:.2f} < {QUALITY_TARGET}, raising detail and retrying…")
            versioned = path.with_name(f"{path.stem}_r{rnd}{path.suffix}")
            path.rename(versioned)
            detail = min(detail + 2, 10)

    tqdm.write(f"\n  Best: {best_path}  (score={best_score:.2f})")
    return best_path


# ── Auto-discovery ────────────────────────────────────────────────────────────

def _find_alignment(slide_path: Path, course_dir: Path) -> Path | None:
    """Find the alignment JSON that was created FOR this slide file.

    Checks the 'slide_file' field inside each alignment JSON for an exact
    or high-similarity match to slide_path.name.  Falls back to filename
    Jaccard similarity if no slide_file field matches.
    """
    align_dir = course_dir / "alignment"
    if not align_dir.exists():
        return None

    slide_name = slide_path.name.lower()
    slide_stem = slide_path.stem.lower()

    # Pass 1: exact match on slide_file field inside alignment JSON
    for f in align_dir.glob("*.json"):
        if f.name.endswith(".compact.json") or f.name.endswith("mapping.json"):
            continue
        try:
            with open(f, encoding="utf-8") as fh:
                data = json.load(fh)
            sf = data.get("slide_file", "")
            if sf and sf.lower() == slide_name:
                return f
        except Exception:
            continue

    # Pass 2: Jaccard similarity on slide_file field
    best_sc, best = 0.0, None
    for f in align_dir.glob("*.json"):
        if f.name.endswith(".compact.json") or f.name.endswith("mapping.json"):
            continue
        try:
            with open(f, encoding="utf-8") as fh:
                data = json.load(fh)
            sf = data.get("slide_file", "")
            if not sf:
                continue
            ta = set(slide_stem.replace("-", " ").split())
            tb = set(Path(sf).stem.lower().replace("-", " ").split())
            if ta and tb:
                sc = len(ta & tb) / len(ta | tb)
                if sc > best_sc:
                    best_sc, best = sc, f
        except Exception:
            continue
    return best if best_sc > 0.3 else None


def _discover_screenshare_lectures(course_dir: Path) -> list[LectureData]:
    """Discover lectures from screen share frame directories.

    Checks the manifest to find screen share (SS) videos, then looks for
    extracted frames in <course_dir>/frames/<video_stem>/.
    """
    frames_dir = course_dir / "frames"
    if not frames_dir.exists():
        return []

    lectures: list[LectureData] = []
    frame_subdirs = sorted([d for d in frames_dir.iterdir() if d.is_dir()])
    for i, fdir in enumerate(frame_subdirs, start=1):
        frame_pngs = sorted(fdir.glob("frame_*.png"))
        if not frame_pngs:
            continue
        # Check for alignment
        align_path = course_dir / "alignment" / f"{fdir.name}.json"
        align = align_path if align_path.exists() else None
        # Verify alignment source is screenshare — skip if not
        if align:
            try:
                with open(align, encoding="utf-8") as f:
                    data = json.load(f)
                if data.get("source") != "screenshare":
                    continue  # This alignment is from traditional slide matching
            except Exception:
                continue  # Corrupted alignment — skip, don't include
        else:
            continue  # No alignment yet — need frame_extractor to run first

        lectures.append(LectureData(
            num=i,
            slide_path=fdir,   # frame_dir serves as slide_path for compatibility
            alignment_path=align,
            file_idx=1,
            source="screenshare",
            frame_dir=fdir,
        ))
    return lectures


def _discover_video_lectures(course_dir: Path,
                             image_source: str = "frames") -> list[LectureData]:
    """Video-centric discovery: iterate captions and pair each with its
    matching slides/frames.

    image_source controls which image source the note generator prefers:
      "frames" (default) — use extracted video frames when available,
                           otherwise fall back to slide PDFs, otherwise
                           transcript-only pseudo-slides.
      "slides"           — prefer slide PDFs, fall back to frames, then
                           transcript-only.

    This is the default discovery mode because each downloaded video should
    produce one note — even when no slide file matches (the note is then
    generated from the transcript alone). The returned LectureData list
    preserves caption order, and `num` is assigned sequentially.
    """
    if image_source not in ("frames", "slides"):
        image_source = "frames"
    captions_dir = course_dir / "captions"
    align_dir    = course_dir / "alignment"
    frames_dir   = course_dir / "frames"
    mat_dir      = course_dir / "materials"

    if not captions_dir.exists():
        return []

    captions = sorted(captions_dir.glob("*.json"))
    if not captions:
        return []

    # Index alignment files by the caption stem they cover. Keys are caption
    # stems; values are (source, alignment_path) so we can prefer screenshare
    # over slide alignments for the same caption (per user memory: frames
    # score higher than PDF slides).
    alignment_by_caption: dict[str, list[Path]] = {}
    screenshare_by_caption: dict[str, Path]    = {}
    if align_dir.exists():
        for af in sorted(align_dir.glob("*.json")):
            if af.name.endswith(".compact.json") or af.name.endswith("mapping.json"):
                continue
            try:
                with open(af, encoding="utf-8") as fh:
                    data = json.load(fh)
            except Exception:
                continue
            # "lecture" field holds the caption stem for both screenshare and
            # slide alignments; fall back to af.stem for old screenshare files.
            cap_stem = data.get("lecture", "") or af.stem
            if data.get("source") == "screenshare":
                # Pick any one screenshare alignment per caption; they are
                # functionally equivalent.
                screenshare_by_caption.setdefault(cap_stem, af)
            elif cap_stem:
                alignment_by_caption.setdefault(cap_stem, []).append(af)

    lectures: list[LectureData] = []
    def _try_frames(num, cap) -> LectureData | None:
        """Return a screenshare LectureData when frames + alignment exist."""
        stem = cap.stem
        if stem not in screenshare_by_caption:
            return None
        frame_dir = frames_dir / stem
        if not (frame_dir.exists() and any(frame_dir.glob("frame_*.png"))):
            return None
        return LectureData(num, frame_dir, screenshare_by_caption[stem],
                           source="screenshare", frame_dir=frame_dir)

    def _try_slides(num, cap) -> list[LectureData]:
        """Return slide-based LectureData(s) for this caption — empty when no
        slide file is available."""
        aligns = alignment_by_caption.get(cap.stem, [])
        out: list[LectureData] = []
        for fi, af in enumerate(aligns, start=1):
            try:
                with open(af, encoding="utf-8") as fh:
                    data = json.load(fh)
            except Exception:
                data = {}
            slide_file = data.get("slide_file", "")
            slide_path = mat_dir / slide_file if slide_file else None
            if slide_path and not slide_path.exists():
                cands = list(mat_dir.rglob(slide_file)) if slide_file else []
                slide_path = cands[0] if cands else None
            if slide_path is None:
                continue
            out.append(LectureData(num, slide_path, af, file_idx=fi))
        return out

    for num, cap in enumerate(captions, start=1):
        tried_first_source = image_source
        first_ld: LectureData | None = None
        first_list: list[LectureData] = []

        if image_source == "slides":
            first_list = _try_slides(num, cap)
            if not first_list:
                # Fall back to frames when the user asked for slides but none
                # are on disk.
                fb = _try_frames(num, cap)
                if fb is not None:
                    tqdm.write(f"  [{cap.stem}] slides requested but none found — "
                               f"falling back to video frames")
                    lectures.append(fb)
                    continue
        else:  # "frames" (default)
            first_ld = _try_frames(num, cap)
            if first_ld is not None:
                lectures.append(first_ld)
                continue
            # Fall back to slides when frames were requested but none extracted.
            first_list = _try_slides(num, cap)
            if first_list:
                tqdm.write(f"  [{cap.stem}] video frames not available — "
                           f"using slide PDFs instead")

        if first_list:
            lectures.extend(first_list)
            continue

        # No alignment + no screenshare frames: emit a transcript-only
        # LectureData. The "slide_path" points to the caption so the note
        # gets a stable identity; generate_section will use the slide_only
        # prompt but with transcript content injected via compact alignment.
        lectures.append(LectureData(num, cap, cap if cap.exists() else None,
                                    source="transcript_only"))
    return lectures


def _discover_lectures(course_dir: Path) -> list[LectureData]:
    exts    = {".pdf", ".pptx", ".ppt", ".docx", ".doc"}

    # ── Check for screen share lectures from extracted frames ────────────────
    ss_lectures = _discover_screenshare_lectures(course_dir)
    if ss_lectures:
        tqdm.write(f"  Found {len(ss_lectures)} screen share lecture(s) from extracted frames")
        # Also discover slide-based lectures for camera recordings
        # (a course may have both screen and camera videos)

    # ── Traditional slide-based discovery ─────────────────────────────────────

    # Common lecture-slide subfolder names, checked in priority order.
    _LECTURE_SUBDIRS = [
        "LectureNotes", "Lecture Slides", "Lecture Notes",
        "Lectures", "Slides", "lecture_notes", "lecture_slides",
    ]

    mat_dir = course_dir / "materials"
    if not mat_dir.exists():
        # Soft return — slide discovery is one of two pathways. The other
        # (video frames / transcripts) is handled by _discover_video_lectures.
        # Letting the caller decide how to surface "nothing to generate"
        # avoids forcing users to download slides when they only want to
        # generate notes from screen-share recordings.
        return []
    lecture_subdir: Path | None = None
    for name in _LECTURE_SUBDIRS:
        candidate = mat_dir / name
        if candidate.exists():
            lecture_subdir = candidate
            break

    if lecture_subdir is not None:
        slide_files = sorted([
            p for p in lecture_subdir.rglob("*")
            if p.suffix.lower() in exts and "image_cache" not in p.name
        ])
    else:
        # Recursively search all of materials/ for lecture-like files.
        # Include a file if it has a lecture number in its name OR lives
        # in a folder whose name suggests lectures.
        _LECTURE_FOLDER_PAT = re.compile(
            r"lecture|slides|notes", re.IGNORECASE)
        _LECTURE_FILE_PAT = re.compile(
            r"(?<![a-zA-Z])[Ll](?:ec(?:ture)?)?[-_ ]?0*\d+", re.IGNORECASE)

        all_candidates = sorted([
            p for p in mat_dir.rglob("*")
            if p.is_file() and p.suffix.lower() in exts
            and "image_cache" not in p.name
        ])

        slide_files = []
        for p in all_candidates:
            # Always include files at the root of materials/
            if p.parent == mat_dir:
                slide_files.append(p)
                continue
            # Include if filename looks like a lecture
            if _LECTURE_FILE_PAT.search(p.stem):
                slide_files.append(p)
                continue
            # Include if any parent folder name suggests lectures
            rel_parts = p.relative_to(mat_dir).parts[:-1]  # folder names
            if any(_LECTURE_FOLDER_PAT.search(part) for part in rel_parts):
                slide_files.append(p)
                continue

    # Group files by lecture number; files without a number get a unique high index
    from collections import defaultdict
    by_num: dict[int, list[Path]] = defaultdict(list)
    auto_num = 1000
    for sp in slide_files:
        m = re.search(r"(?<![a-zA-Z])[Ll](?:ec(?:ture)?)?[-_ ]?0*(\d+)", sp.stem)
        if m:
            lec_num = int(m.group(1))
        else:
            lec_num  = auto_num
            auto_num += 1
        by_num[lec_num].append(sp)

    lectures: list[LectureData] = []
    for lec_num in sorted(by_num.keys()):
        files = sorted(by_num[lec_num])   # alphabetical within same lecture number
        for file_idx, sp in enumerate(files, start=1):
            align = _find_alignment(sp, course_dir)
            lectures.append(LectureData(lec_num, sp, align, file_idx=file_idx))

    # Deduplicate within each lecture number: when multiple slide files for
    # the same lecture share the same alignment file, keep only the best one
    # (prefer main slides over pre-lecture/review/annotated copies).
    # This prevents duplicate content from e.g. "Lecture 6.pdf" and
    # "Lecture 6 Pre Lecture.pdf" both aligning to the same transcript.
    _SECONDARY_PAT = re.compile(
        r"pre[_\s-]?lec|review|with\s*notes|ann\d|annotated",
        re.IGNORECASE)
    deduped: list[LectureData] = []
    from itertools import groupby as _gby
    for lec_num, grp in _gby(sorted(lectures, key=lambda x: x.num),
                              key=lambda x: x.num):
        group = list(grp)
        if len(group) <= 1:
            deduped.extend(group)
            continue
        # Within this lecture number, find files sharing the same alignment
        seen: dict[str, LectureData] = {}
        for ld in group:
            if not ld.alignment_path:
                deduped.append(ld)
                continue
            key = str(ld.alignment_path)
            if key not in seen:
                seen[key] = ld
            else:
                existing = seen[key]
                # Prefer the file that is NOT a secondary variant
                if _SECONDARY_PAT.search(existing.slide_path.stem) and \
                   not _SECONDARY_PAT.search(ld.slide_path.stem):
                    seen[key] = ld  # replace with main file
                # else: keep existing
        deduped.extend(seen.values())

    if len(deduped) < len(lectures):
        tqdm.write(f"  Deduplicated: {len(lectures)} → {len(deduped)} lecture(s) "
                   f"(removed slide variants sharing the same alignment)")
    lectures = deduped

    # Renumber file_idx within each lecture number after dedup
    final: list[LectureData] = []
    for _, grp in _gby(sorted(lectures, key=lambda x: x.num), key=lambda x: x.num):
        for fi, ld in enumerate(grp, 1):
            ld.file_idx = fi
            final.append(ld)
    lectures = final

    # Combine with screen share lectures: screenshare replaces slide-based
    # entries that cover the same lecture content.
    if ss_lectures:
        # Try to match each screenshare to a slide-based lecture by
        # extracting a lecture number from the video/frame directory name.
        _lec_num_pat = re.compile(
            r"(?<![a-zA-Z])[Ll](?:ec(?:ture)?)?[-_ ]?0*(\d+)", re.IGNORECASE)
        replaced_nums: set[int] = set()
        for ss in ss_lectures:
            m = _lec_num_pat.search(ss.slide_path.name)
            if m:
                target_num = int(m.group(1))
                # Find and replace the slide-based entry with this number
                for i, ld in enumerate(lectures):
                    if ld.num == target_num and ld.source == "slides":
                        ss.num = target_num
                        lectures[i] = ss
                        replaced_nums.add(target_num)
                        break
                else:
                    # No matching slide lecture — append with a new number
                    max_num = max((ld.num for ld in lectures), default=0)
                    ss.num = max_num + 1
                    lectures.append(ss)
            else:
                max_num = max((ld.num for ld in lectures), default=0)
                ss.num = max_num + 1
                lectures.append(ss)

        if replaced_nums:
            tqdm.write(f"  Screenshare replaces slide-based for lecture(s) "
                       f"{sorted(replaced_nums)}")

    # ── Final pass: add any aligned captions not yet covered ─────────────
    # Some videos have alignments but don't match any slide file or
    # screenshare frame directory.  Add them so every video gets a note.
    covered_aligns = {str(ld.alignment_path) for ld in lectures if ld.alignment_path}
    align_dir = course_dir / "alignment"
    if align_dir.exists():
        max_num = max((ld.num for ld in lectures), default=0)
        added = 0
        for af in sorted(align_dir.glob("*.json")):
            if af.name.endswith(".compact.json") or af.name.endswith("mapping.json"):
                continue
            if str(af) in covered_aligns:
                continue
            try:
                with open(af, encoding="utf-8") as fh:
                    data = json.load(fh)
                if data.get("source") == "screenshare":
                    # Screenshare without frame dir — skip (needs frame_extractor)
                    frame_dir = course_dir / "frames" / af.stem
                    if not frame_dir.exists():
                        continue
                    max_num += 1
                    lectures.append(LectureData(
                        max_num, frame_dir, af, source="screenshare", frame_dir=frame_dir))
                else:
                    slide_file = data.get("slide_file", "")
                    slide_path = course_dir / "materials" / slide_file if slide_file else af
                    if not slide_path.exists():
                        # Try finding in subdirs
                        candidates = list((course_dir / "materials").rglob(slide_file)) if slide_file else []
                        slide_path = candidates[0] if candidates else af
                    max_num += 1
                    lectures.append(LectureData(max_num, slide_path, af))
                added += 1
            except Exception:
                continue
        if added:
            tqdm.write(f"  Added {added} video(s) without matching slides")

    return lectures


# ── CLI ───────────────────────────────────────────────────────────────────────

def main() -> None:
    global NOTE_LANGUAGE, SHOW_SCORE, NOTE_MODEL
    parser = argparse.ArgumentParser(description="Generate course lecture notes")
    parser.add_argument("--course",      metavar="ID")
    parser.add_argument("--slides",      metavar="PATH")
    parser.add_argument("--alignment",   metavar="PATH")
    parser.add_argument("--lecture-num", metavar="N", type=int, default=1)
    parser.add_argument("--course-name", metavar="NAME", default="")
    parser.add_argument("--out",         metavar="PATH")
    parser.add_argument("--detail",      metavar="0-10", type=int, default=DETAIL_LEVEL)
    parser.add_argument("--format",      metavar="md|mdx", default=OUTPUT_FORMAT,
                        choices=["md", "mdx"])
    parser.add_argument("--iterate",     action="store_true")
    parser.add_argument("--merge-only",  action="store_true",
                        help="Skip generation; just merge existing sections into final note")
    parser.add_argument("--force",       action="store_true",
                        help="Re-generate sections even if they already exist")
    parser.add_argument("--lectures",    metavar="N-N or N,N,N", default="",
                        help="Filter lectures, e.g. '1-5' or '1,2,3'")
    parser.add_argument("--per-video",  action="store_true",
                        help="[deprecated: per-video is now the default] "
                             "Generate one note file per video/lecture.")
    parser.add_argument("--merged",     action="store_true",
                        help="Generate one combined course-wide note file "
                             "instead of the default one-note-per-video layout.")
    parser.add_argument("--slide-centric", action="store_true",
                        help="[legacy] Iterate slide files first when building "
                             "the lecture list; by default the pipeline is "
                             "video-centric and iterates captions.")
    parser.add_argument("--image-source", metavar="SRC", default="frames",
                        choices=["frames", "slides"],
                        help="Image source for notes: 'frames' (video "
                             "screenshots, default) or 'slides' (PDF slide "
                             "renders). Falls back to the other source when "
                             "the preferred one is not available.")
    parser.add_argument("--language",   metavar="LANG", default=None,
                        choices=["en", "zh"],
                        help="Note language: en (English) or zh (Chinese). "
                             "Overrides the NOTE_LANGUAGE constant.")
    parser.add_argument("--score",      action="store_true",
                        help="Enable self-scoring (dev mode: prints coverage/terminology scores)")
    parser.add_argument("--model",      metavar="MODEL", default=None,
                        help="Override NOTE_MODEL (e.g. codex-cli, claude-cli, gemini-2.5-flash)")
    args = parser.parse_args()

    if args.score:
        SHOW_SCORE = True

    if args.model:
        global NOTE_MODEL
        NOTE_MODEL = args.model
        print(f"Note model: {NOTE_MODEL}")

    if args.language:
        NOTE_LANGUAGE = args.language
        print(f"Note language: {NOTE_LANGUAGE}")

    if args.course:
        course_dir  = COURSE_DATA_DIR / args.course
        course_name = args.course_name or f"CS{args.course}"
        if args.slide_centric:
            lectures = _discover_lectures(course_dir)
        else:
            lectures = _discover_video_lectures(course_dir,
                                                image_source=args.image_source)
            if not lectures:
                # No captions yet — fall back to slide-based discovery so users
                # can still generate notes from raw slides without running
                # the transcribe step.
                print("  No captions found — falling back to slide-based discovery")
                lectures = _discover_lectures(course_dir)
        if args.lectures:
            sel: set[int] = set()
            for part in args.lectures.split(","):
                part = part.strip()
                if "-" in part:
                    a, b = part.split("-", 1)
                    sel.update(range(int(a), int(b) + 1))
                elif part.isdigit():
                    sel.add(int(part))
            lectures = [l for l in lectures if l.num in sel]
        if not lectures:
            captions_dir = course_dir / "captions"
            mat_dir      = course_dir / "materials"
            print(f"[error] Nothing to generate notes from under {course_dir}")
            if not captions_dir.exists() and not mat_dir.exists():
                print("  Neither captions/ nor materials/ exists. Download "
                      "videos and transcribe them, or download slide PDFs "
                      "(or both) and try again.")
            elif not captions_dir.exists():
                print("  No transcripts in captions/. Transcribe videos "
                      "first, or rerun with --image-source slides if you "
                      "only want to use the slide PDFs.")
            elif not mat_dir.exists():
                print("  No slides in materials/. Run with --image-source "
                      "frames (default) so the screenshare path is used, "
                      "or download lecture slides and rerun.")
            else:
                print("  Both captions/ and materials/ exist but no "
                      "lecture pairs were discovered. Check filenames + "
                      "alignment output.")
            sys.exit(1)

        print(f"Found {len(lectures)} lectures:")
        for ld in lectures:
            src = f" [{ld.source}]" if ld.source != "slides" else ""
            a  = "+ alignment" if ld.alignment_path else "(slide-only)"
            fi = f" [part {ld.file_idx}]" if ld.file_idx > 1 else ""
            print(f"  L{ld.num}{fi}{src}: {ld.slide_path.name}  {a}")

        ext_out      = ".mdx" if args.format == "mdx" else ".md"
        out_path     = Path(args.out) if args.out else \
                       course_dir / "notes" / f"{course_name}_notes{ext_out}"
        sections_dir = out_path.parent / "sections"

        # Auto-force when language changed from the previously cached language
        lang_marker = sections_dir / ".language"
        if not args.force and lang_marker.exists():
            cached_lang = lang_marker.read_text(encoding="utf-8").strip()
            if cached_lang != NOTE_LANGUAGE:
                args.force = True
                print(f"  Language changed ({cached_lang} → {NOTE_LANGUAGE}), "
                      f"force-regenerating sections.")

        # Per-video is the default; pass --merged to produce a single combined
        # note file instead.
        if not args.merged:
            out_dir = out_path.parent
            generate_per_video_notes(
                course_name, lectures, out_dir,
                detail=args.detail, fmt=args.format, force=args.force,
            )
            return

        if args.merge_only:
            # Load lectures so titles are available
            out_path.parent.mkdir(parents=True, exist_ok=True)
            for ld in lectures:
                ld.load(out_path.parent)
            all_slides  = [s for ld in lectures for s in ld.slides]
            all_compact = [c for ld in lectures for c in ld.compact_slides]
            merge_sections(course_name, lectures, sections_dir, out_path,
                           all_slides, all_compact)
            return

        if args.iterate:
            generate_with_iteration(course_name, lectures, out_path, fmt=args.format)
        else:
            generate_course_notes(course_name, lectures, out_path,
                                  detail=args.detail, fmt=args.format,
                                  force=args.force)
        return

    if not args.slides:
        parser.error("Provide --course or --slides")

    sp = Path(args.slides)
    if not sp.exists():
        print(f"[error] Not found: {sp}"); sys.exit(1)

    align    = Path(args.alignment) if args.alignment else None
    ld       = LectureData(args.lecture_num, sp, align)
    name     = args.course_name or sp.parent.parent.parent.name
    ext_out  = ".mdx" if args.format == "mdx" else ".md"
    out_path = Path(args.out) if args.out else \
               sp.parent.parent.parent / "notes" / f"{name}_notes{ext_out}"

    fn = generate_with_iteration if args.iterate else \
         lambda *a, **kw: generate_course_notes(*a, **kw)
    fn(name, [ld], out_path, fmt=args.format,
       **({} if args.iterate else {"detail": args.detail}))


if __name__ == "__main__":
    try:
        main()
    except KeyboardInterrupt:
        print("\n[info] Interrupted by user.")
        sys.exit(0)
    except SystemExit:
        raise
    except Exception as _exc:
        import traceback
        print(f"\n[error] Unexpected error: {_exc}")
        traceback.print_exc()
        sys.exit(1)
